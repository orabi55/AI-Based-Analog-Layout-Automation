import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Color Constants
ACADEMIC_BLUE = RGBColor(11, 37, 69)       # #0B2545
WARM_GOLD = RGBColor(197, 160, 89)        # #C5A059
SOFT_CREAM = RGBColor(253, 251, 247)       # #FDFBF7
TEXT_CHARCOAL = RGBColor(43, 43, 43)       # #2B2B2B
ROYAL_BLUE = RGBColor(26, 75, 132)         # #1A4B84
CODE_GREEN = RGBColor(28, 124, 84)         # #1C7C54
CODE_BLUE = RGBColor(0, 95, 115)           # #005F73
LIGHT_BLUE_BG = RGBColor(240, 244, 248)    # #F0F4F8
WHITE = RGBColor(255, 255, 255)
LIGHT_GREY = RGBColor(220, 224, 230)

def get_image_path(filename):
    """Robust image path resolution."""
    paths = [
        os.path.join("..", "Thesis", "figures", filename),
        os.path.join("Thesis", "figures", filename),
        os.path.join("figures", filename),
        filename
    ]
    for p in paths:
        if os.path.exists(p):
            return p
    return None

def add_solid_background(slide, color):
    """Draws a full-slide rectangle as background."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect

def add_shadow_card(slide, left, top, width, height, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, color=LIGHT_GREY):
    """Draws a flat offset shape as a drop shadow to create a premium 3D layered card effect."""
    offset = Inches(0.05)
    if width > Inches(3.0) or height > Inches(3.0):
        offset = Inches(0.08)
    shadow = slide.shapes.add_shape(shape_type, left + offset, top + offset, width, height)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = color
    shadow.line.fill.background()
    return shadow

def add_slide_header(slide, title_text):
    """Creates a top banner with Academic Blue background and Gold line."""
    # Banner
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.9))
    banner.fill.solid()
    banner.fill.fore_color.rgb = ACADEMIC_BLUE
    banner.line.fill.background()
    
    # Title text box
    tx_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.7))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Calibri"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    
    # Gold accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.9), Inches(13.333), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = WARM_GOLD
    line.line.fill.background()

def add_slide_footer(slide, current_page=None, total_pages=None, is_dark=False):
    """Adds small grey/white footer text with page numbers and a horizontal separator line."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls
    
    # Text and line color based on slide background
    text_color = RGBColor(200, 200, 200) if is_dark else RGBColor(128, 128, 128)
    line_color = RGBColor(100, 100, 100) if is_dark else RGBColor(230, 230, 230)
    
    # Horizontal separator line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(7.0), Inches(12.533), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = line_color
    line.line.fill.background()
    
    # Left footer text box (descriptive text)
    tx_box_left = slide.shapes.add_textbox(Inches(0.4), Inches(7.1), Inches(8.0), Inches(0.3))
    tf_left = tx_box_left.text_frame
    tf_left.margin_left = Inches(0)
    tf_left.margin_right = Inches(0)
    tf_left.margin_top = Inches(0)
    tf_left.margin_bottom = Inches(0)
    p_left = tf_left.paragraphs[0]
    p_left.text = "AI-Based Analog Layout Automation | Ain Shams University & SI-Vision Collaboration"
    p_left.font.name = "Calibri"
    p_left.font.size = Pt(10)
    p_left.font.color.rgb = text_color
    p_left.alignment = PP_ALIGN.LEFT
    
    # Right footer text box (dynamic slidenum field)
    tx_box_right = slide.shapes.add_textbox(Inches(8.8), Inches(7.1), Inches(4.133), Inches(0.3))
    tf_right = tx_box_right.text_frame
    tf_right.margin_left = Inches(0)
    tf_right.margin_right = Inches(0)
    tf_right.margin_top = Inches(0)
    tf_right.margin_bottom = Inches(0)
    p_right = tf_right.paragraphs[0]
    p_right.text = "Slide "
    p_right.alignment = PP_ALIGN.RIGHT
    
    # XML color hex based on dark or light background
    color_hex = "C8C8C8" if is_dark else "808080"
    
    fld_xml = (
        '<a:fld %s id="{1F4E2DE4-8ADA-4D4E-9951-90A1D26586E7}" type="slidenum">'
        '<a:rPr lang="en-US" smtClean="0" sz="1000">'
        '<a:latin typeface="Calibri"/>'
        '<a:solidFill>'
        '<a:srgbClr val="%s"/>'
        '</a:solidFill>'
        '</a:rPr>'
        '<a:t>#</a:t>'
        '</a:fld>' % (nsdecls('a'), color_hex)
    )
    p_right._p.append(parse_xml(fld_xml))
    
    p_right.font.name = "Calibri"
    p_right.font.size = Pt(10)
    p_right.font.color.rgb = text_color



def add_content_box(slide, left, top, width, height, title, bullets):
    """Adds a rounded rectangle content box with bullet points."""
    # Offset shadow card
    add_shadow_card(slide, left, top, width, height, MSO_SHAPE.ROUNDED_RECTANGLE, LIGHT_GREY)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BLUE_BG
    box.line.color.rgb = WARM_GOLD
    box.line.width = Pt(1)
    
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_right = Inches(0.2)
    
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = "Calibri"
    p_title.font.size = Pt(16)
    p_title.font.bold = True
    p_title.font.color.rgb = ACADEMIC_BLUE
    p_title.space_after = Pt(8)
    
    for b in bullets:
        p = tf.add_paragraph()
        p.text = f"• {b}"
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_CHARCOAL
        p.space_after = Pt(4)
        p.level = 0

def add_code_box(slide, left, top, width, height, title, code_lines):
    """Adds a dark premium IDE-style box for displaying code snippets with basic syntax highlighting."""
    # Gold drop glow shadow
    add_shadow_card(slide, left, top, width, height, MSO_SHAPE.RECTANGLE, ACADEMIC_BLUE)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(28, 28, 30) # Dark grey background
    box.line.color.rgb = WARM_GOLD # Gold border
    box.line.width = Pt(1.5)
    
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_right = Inches(0.2)
    
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = "Calibri"
    p_title.font.size = Pt(14)
    p_title.font.bold = True
    p_title.font.color.rgb = WARM_GOLD
    p_title.space_after = Pt(8)
    
    # Syntax highlighting color definitions (OneDark theme inspired)
    COLOR_KEYWORD = RGBColor(224, 108, 117)  # Soft red (#E06C75)
    COLOR_STRING = RGBColor(152, 195, 121)   # Soft green (#98C379)
    COLOR_NUMBER = RGBColor(209, 154, 102)   # Soft orange (#D19A66)
    COLOR_COMMENT = RGBColor(92, 99, 112)    # Muted gray (#5C6370)
    COLOR_DEFAULT = RGBColor(171, 178, 191)  # Off-white (#ABB2BF)
    COLOR_BRACKET = RGBColor(86, 182, 194)   # Cyan (#56B6C2)
    
    # Simple regexes to tokenize code lines
    # We match: comments, strings, keywords, numbers, brackets, and whitespace/identifiers
    token_pattern = re.compile(
        r'(?P<comment>#.*|//.*)'
        r'|(?P<string>"[^"]*")'
        r'|(?P<keyword>\b(?:if|else|foreach|set|read|open|close|proc|return|class|def|dict|nodes|edges|id|type|nets)\b)'
        r'|(?P<number>\b\d+(?:\.\d+)?(?:e-?\d+)?\b)'
        r'|(?P<bracket>[{}()\[\]])'
        r'|(?P<other>[^\s{}()\[\]"#]+|\s+)'
    )
    
    for line in code_lines:
        p = tf.add_paragraph()
        p.space_after = Pt(1)
        p.font.name = "Consolas"
        p.font.size = Pt(9.5)
        
        # Check if line is completely a comment
        trimmed = line.strip()
        if trimmed.startswith("#") or trimmed.startswith("//"):
            run = p.add_run()
            run.text = line
            run.font.color.rgb = COLOR_COMMENT
            continue
            
        matches = list(token_pattern.finditer(line))
        if not matches:
            run = p.add_run()
            run.text = line
            run.font.color.rgb = COLOR_DEFAULT
            continue
            
        for match in matches:
            text = match.group(0)
            kind = match.lastgroup
            
            run = p.add_run()
            run.text = text
            
            if kind == "comment":
                run.font.color.rgb = COLOR_COMMENT
            elif kind == "string":
                run.font.color.rgb = COLOR_STRING
            elif kind == "keyword":
                run.font.color.rgb = COLOR_KEYWORD
                run.font.bold = True
            elif kind == "number":
                run.font.color.rgb = COLOR_NUMBER
            elif kind == "bracket":
                run.font.color.rgb = COLOR_BRACKET
            else:
                run.font.color.rgb = COLOR_DEFAULT

def add_table_box(slide, left, top, width, height, headers, rows_data):
    """Adds a styled native PowerPoint table."""
    num_rows = len(rows_data) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table
    
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACADEMIC_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.name = "Calibri"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
    for row_idx, row in enumerate(rows_data, 1):
        for col_idx, val in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(val)
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = LIGHT_BLUE_BG
            else:
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.name = "Calibri"
                p.font.size = Pt(10)
                p.font.color.rgb = TEXT_CHARCOAL
                p.alignment = PP_ALIGN.CENTER

def add_boxed_image(slide, left, top, width, height, filename, label):
    """Draws a boxed figure (equivalent to LaTeX fbox) with a text caption and drop shadow."""
    # Offset shadow card
    add_shadow_card(slide, left, top, width, height, MSO_SHAPE.RECTANGLE, LIGHT_GREY)
    
    # Draw border shape
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    border.fill.solid()
    border.fill.fore_color.rgb = WHITE
    border.line.color.rgb = ACADEMIC_BLUE
    border.line.width = Pt(1.5)
    
    # Try inserting picture
    img_path = get_image_path(filename)
    inset = Inches(0.05)
    if img_path:
        slide.shapes.add_picture(img_path, left + inset, top + inset, width - 2*inset, height - 2*inset)
    else:
        # Fallback label inside box
        tx_box = slide.shapes.add_textbox(left, top + Inches(0.5), width, height - Inches(1.0))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"[Image: {filename}]"
        p.font.name = "Calibri"
        p.font.size = Pt(11)
        p.font.color.rgb = WARM_GOLD
        p.alignment = PP_ALIGN.CENTER
        
    # Draw caption below
    label_box = slide.shapes.add_textbox(left, top + height, width, Inches(0.35))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = "Calibri"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = TEXT_CHARCOAL
    p.alignment = PP_ALIGN.CENTER

def add_block_arrow(slide, left, top, width, height, color=WARM_GOLD, orient="R"):
    """Adds a block arrow. orient can be 'R' (Right), 'L' (Left), 'U' (Up), 'D' (Down)."""
    shape_type = MSO_SHAPE.RIGHT_ARROW
    if orient == "L":
        shape_type = MSO_SHAPE.LEFT_ARROW
    elif orient == "U":
        shape_type = MSO_SHAPE.UP_ARROW
    elif orient == "D":
        shape_type = MSO_SHAPE.DOWN_ARROW
        
    arrow = slide.shapes.add_shape(shape_type, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow

def add_cover_slide(slide, data):
    """Stunning title page matching the capsule format of the reference slides."""
    add_solid_background(slide, RGBColor(245, 248, 252))
    
    # Subtitle: GRADUATION PROJECT
    tx_box_sub = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(0.4))
    tf_sub = tx_box_sub.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "GRADUATION PROJECT"
    p_sub.font.name = "Calibri"
    p_sub.font.size = Pt(14)
    p_sub.font.bold = True
    p_sub.font.color.rgb = ROYAL_BLUE
    p_sub.alignment = PP_ALIGN.CENTER
    
    # Title
    tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.1), Inches(11.333), Inches(1.6))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data["title"]
    p.font.name = "Calibri"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACADEMIC_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Student Capsules (6 items in 2 rows of 3)
    students = [
        {"name": "Mohammed Magdy", "role": "SPICE Parser & Ingestion"},
        {"name": "Noureldin Ayman", "role": "LangGraph Placer Agent"},
        {"name": "Ahmed Khairat", "role": "DRC Critic Spatial Healing"},
        {"name": "Mohammed Wael", "role": "PySide6 Canvas GUI"},
        {"name": "Omar Ahmed", "role": "Chatbot Co-Pilot Agent"},
        {"name": "Eman Sherif", "role": "Diffusion Compaction"}
    ]
    
    lefts = [Inches(1.0), Inches(5.0), Inches(9.0)]
    tops = [Inches(3.3), Inches(4.3)]
    
    for i, s in enumerate(students):
        r = i // 3
        c = i % 3
        
        # capsule box
        cap = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lefts[c], tops[r], Inches(3.3), Inches(0.75))
        cap.fill.solid()
        cap.fill.fore_color.rgb = LIGHT_BLUE_BG
        cap.line.color.rgb = WARM_GOLD
        cap.line.width = Pt(1)
        
        tf_cap = cap.text_frame
        tf_cap.word_wrap = True
        tf_cap.margin_left = Inches(0.1)
        tf_cap.margin_top = Inches(0.08)
        
        p_name = tf_cap.paragraphs[0]
        p_name.text = s["name"]
        p_name.font.name = "Calibri"
        p_name.font.size = Pt(13)
        p_name.font.bold = True
        p_name.font.color.rgb = ACADEMIC_BLUE
        p_name.alignment = PP_ALIGN.CENTER
        
        p_role = tf_cap.add_paragraph()
        p_role.text = s["role"]
        p_role.font.name = "Calibri"
        p_role.font.size = Pt(9.5)
        p_role.font.color.rgb = TEXT_CHARCOAL
        p_role.alignment = PP_ALIGN.CENTER
        
    # Submission details
    tx_sub = slide.shapes.add_textbox(Inches(1.0), Inches(5.3), Inches(11.333), Inches(0.8))
    tf_s = tx_sub.text_frame
    p_s = tf_s.paragraphs[0]
    p_s.text = "Submitted To: Prof. DiaaEldin S. Khalil"
    p_s.font.name = "Calibri"
    p_s.font.size = Pt(13)
    p_s.font.bold = True
    p_s.font.color.rgb = TEXT_CHARCOAL
    p_s.alignment = PP_ALIGN.CENTER
    
    p_uni = tf_s.add_paragraph()
    p_uni.text = "SI-Vision  |  Ain Shams University Faculty of Engineering"
    p_uni.font.name = "Calibri"
    p_uni.font.size = Pt(12)
    p_uni.font.color.rgb = ROYAL_BLUE
    p_uni.space_before = Pt(4)
    p_uni.alignment = PP_ALIGN.CENTER
    
    # Logos side-by-side
    logo_aus_path = get_image_path("ASU.png")
    logo_si_path = get_image_path("SI.png")
    
    if logo_aus_path:
        slide.shapes.add_picture(logo_aus_path, Inches(1.2), Inches(6.0), height=Inches(0.9))
    if logo_si_path:
        slide.shapes.add_picture(logo_si_path, Inches(10.5), Inches(6.0), height=Inches(0.8))

def add_revolution_slide(slide, idx, total_pages):
    """Stunning multi-panel slide describing the abstraction flow."""
    add_slide_header(slide, "Revolutionising IC Design: The AI-Assisted Symbolic Layout")
    add_slide_footer(slide, idx, total_pages)
    
    # Left column: Complexity Barrier & Symbolic Abstraction
    # Complexity barrier box
    add_shadow_card(slide, Inches(0.6), Inches(1.4), Inches(5.4), Inches(2.3), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_barrier = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(5.4), Inches(2.3))
    box_barrier.fill.solid()
    box_barrier.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_barrier.line.color.rgb = ROYAL_BLUE
    box_barrier.line.width = Pt(1)
    
    tf_bar = box_barrier.text_frame
    tf_bar.word_wrap = True
    p_bar_title = tf_bar.paragraphs[0]
    p_bar_title.text = "The Complexity Barrier"
    p_bar_title.font.name = "Calibri"
    p_bar_title.font.size = Pt(16)
    p_bar_title.font.bold = True
    p_bar_title.font.color.rgb = ROYAL_BLUE
    
    p_bar_b1 = tf_bar.add_paragraph()
    p_bar_b1.text = "• Manual sub-micron geometries are too complex for direct AI coordinate mapping."
    p_bar_b1.font.name = "Calibri"
    p_bar_b1.font.size = Pt(11.5)
    p_bar_b1.font.color.rgb = TEXT_CHARCOAL
    p_bar_b1.space_before = Pt(4)
    
    p_bar_b2 = tf_bar.add_paragraph()
    p_bar_b2.text = "• Inflexible analytical cost functions fail to capture high-level qualitative designer preferences."
    p_bar_b2.font.name = "Calibri"
    p_bar_b2.font.size = Pt(11.5)
    p_bar_b2.font.color.rgb = TEXT_CHARCOAL
    
    # Symbolic Abstraction box
    add_shadow_card(slide, Inches(0.6), Inches(3.9), Inches(5.4), Inches(2.9), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_abs = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(3.9), Inches(5.4), Inches(2.9))
    box_abs.fill.solid()
    box_abs.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_abs.line.color.rgb = ROYAL_BLUE
    box_abs.line.width = Pt(1)
    
    tf_abs = box_abs.text_frame
    tf_abs.word_wrap = True
    p_abs_title = tf_abs.paragraphs[0]
    p_abs_title.text = "Intelligent Symbolic Representation"
    p_abs_title.font.name = "Calibri"
    p_abs_title.font.size = Pt(16)
    p_abs_title.font.bold = True
    p_abs_title.font.color.rgb = ROYAL_BLUE
    
    p_abs_b1 = tf_abs.add_paragraph()
    p_abs_b1.text = "• Decouples layout: strategic planning (LLM) and geometric grid-snapping compilation."
    p_abs_b1.font.name = "Calibri"
    p_abs_b1.font.size = Pt(11.5)
    p_abs_b1.font.color.rgb = TEXT_CHARCOAL
    p_abs_b1.space_before = Pt(4)
    
    p_abs_b2 = tf_abs.add_paragraph()
    p_abs_b2.text = "• Translates absolute shapes into logical parameters (ordering, orientations, finger splits) readable by both humans and AI."
    p_abs_b2.font.name = "Calibri"
    p_abs_b2.font.size = Pt(11.5)
    p_abs_b2.font.color.rgb = TEXT_CHARCOAL
    
    # Right column: The 3-Step Optimisation Workflow
    tx_wf = slide.shapes.add_textbox(Inches(6.4), Inches(1.2), Inches(6.3), Inches(0.4))
    tx_wf.text_frame.paragraphs[0].text = "The 3-Step Optimisation Workflow"
    tx_wf.text_frame.paragraphs[0].font.name = "Calibri"
    tx_wf.text_frame.paragraphs[0].font.size = Pt(17)
    tx_wf.text_frame.paragraphs[0].font.bold = True
    tx_wf.text_frame.paragraphs[0].font.color.rgb = ACADEMIC_BLUE
    
    step_y = [Inches(1.6), Inches(2.8), Inches(4.0)]
    step_titles = [
        "1. Ingestion & AI Proposal",
        "2. Human-in-the-Loop Refinement",
        "3. Snapping & Physical Conversion"
    ]
    step_desc = [
        "Parses SPICE netlists, detects symmetry matching groups, and suggests initial relative configurations.",
        "Designer uses an interactive canvas editor to modify placements and type conversational chat prompt overrides.",
        "Deterministic compiler snapping compresses layout boundaries and outputs binary GDSII/OASIS streams."
    ]
    
    for s_idx in range(3):
        add_shadow_card(slide, Inches(6.4), step_y[s_idx], Inches(6.3), Inches(1.1), MSO_SHAPE.ROUNDED_RECTANGLE)
        sbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.4), step_y[s_idx], Inches(6.3), Inches(1.1))
        sbox.fill.solid()
        sbox.fill.fore_color.rgb = SOFT_CREAM
        sbox.line.color.rgb = WARM_GOLD
        sbox.line.width = Pt(1)
        
        tf_s = sbox.text_frame
        tf_s.word_wrap = True
        tf_s.margin_left = Inches(0.15)
        tf_s.margin_top = Inches(0.08)
        
        p_st = tf_s.paragraphs[0]
        p_st.text = step_titles[s_idx]
        p_st.font.name = "Calibri"
        p_st.font.size = Pt(12)
        p_st.font.bold = True
        p_st.font.color.rgb = ACADEMIC_BLUE
        
        p_sd = tf_s.add_paragraph()
        p_sd.text = step_desc[s_idx]
        p_sd.font.name = "Calibri"
        p_sd.font.size = Pt(10)
        p_sd.font.color.rgb = TEXT_CHARCOAL
        p_sd.space_before = Pt(3)
        
    # Table of System components functions
    headers = ["System Component", "Primary Function"]
    rows = [
        ["EDA Interface", "Parses netlist and translates symbolic parameters back to database layout."],
        ["AI / LLM Engine", "Analyzes block relationships and optimizes layouts via conversational loops."],
        ["Visualizer GUI", "Provides an interactive canvas editor for manual overrides and chat inputs."]
    ]
    add_table_box(slide, Inches(6.4), Inches(5.3), Inches(6.3), Inches(1.5), headers, rows)

def add_timeline_slide(slide, idx, total_pages):
    """Draws a professional vertical timeline slide."""
    add_slide_header(slide, "System Execution Workflow")
    add_slide_footer(slide, idx, total_pages)
    
    # Vertical line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.5), Inches(0.05), Inches(5.2))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_GREY
    line.line.fill.background()
    
    steps = [
        {"title": "Import & Ingest", "desc": "Ingests SPICE netlist and resolves PDK fin/poly track snapping bounds.", "color": ROYAL_BLUE},
        {"title": "Symmetry Match", "desc": "NetworkX bipartite graph detects matching pairs (diff pairs, mirrors).", "color": CODE_GREEN},
        {"title": "AI Initial Placement", "desc": "LangGraph Specialist agent generates initial symbolic layouts.", "color": CODE_BLUE},
        {"title": "DRC Self-Healing", "desc": "Sweepline Spacing Critic loops corrections to Planner node.", "color": WARM_GOLD},
        {"title": "Visual Review & Chat", "desc": "Designer inspects layout canvas and overrides parameters via chat.", "color": RGBColor(230, 90, 40)},
        {"title": "Physical Exporter", "desc": "Gdstk generates binary GDSII/OASIS; watchdog inserts to database.", "color": ACADEMIC_BLUE}
    ]
    
    y_starts = [1.5, 2.35, 3.2, 4.05, 4.9, 5.75]
    
    for s_idx, step in enumerate(steps):
        y_pos = Inches(y_starts[s_idx])
        
        # Circle node on line
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.375), y_pos + Inches(0.08), Inches(0.3), Inches(0.3))
        circle.fill.solid()
        circle.fill.fore_color.rgb = step["color"]
        circle.line.color.rgb = WHITE
        circle.line.width = Pt(1.5)
        
        # Text box
        tx = slide.shapes.add_textbox(Inches(1.8), y_pos, Inches(10.5), Inches(0.75))
        tf = tx.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        
        p_t = tf.paragraphs[0]
        p_t.text = step["title"]
        p_t.font.name = "Calibri"
        p_t.font.size = Pt(13.5)
        p_t.font.bold = True
        p_t.font.color.rgb = step["color"]
        
        p_d = tf.add_paragraph()
        p_d.text = step["desc"]
        p_d.font.name = "Calibri"
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_CHARCOAL
        p_d.space_before = Pt(1)

def add_four_panel_slide(slide, idx, total_pages):
    """Draws a professional 4-panel vertical card layout."""
    add_slide_header(slide, "System Architecture: 4-Column Panel Grid")
    add_slide_footer(slide, idx, total_pages)
    
    panel_lefts = [0.8, 3.8, 6.8, 9.8]
    panel_width = 2.7
    panel_height = 5.3
    
    panels_data = [
        {
            "title": "EDA Ingestion",
            "frame_color": ACADEMIC_BLUE,
            "bullets": [
                "Netlist Ingestion: Ingests flat SPICE netlist files (.sp, .cdl).",
                "Device Ingest: Resolves widths (W), fingers (M), and lengths (L).",
                "PDK Snapping Lookups: Maps poly tracks and active fin tracks.",
                "Symmetry Matching: Identifies differential pairs and current mirrors."
            ]
        },
        {
            "title": "AI/LLM Engine",
            "frame_color": ROYAL_BLUE,
            "bullets": [
                "LangGraph State Machine: Coordinates multi-agent layout crew.",
                "Builder & Specialist: Plans relative row configurations.",
                "DRC Spacing Critic: Checks physical overlaps via Sweepline.",
                "ReAct Loop: Iterates coordinate correction constraints."
            ]
        },
        {
            "title": "Human Co-Pilot",
            "frame_color": CODE_BLUE,
            "bullets": [
                "Conversational Panel: Chat interface in PySide6 GUI.",
                "Intent Routing: Routes prompts (move, swap, align, abut).",
                "QThread Worker: Background thread prevents GUI canvas lag.",
                "Command XML Queue: Batch updates via QTimer."
            ]
        },
        {
            "title": "Outputs",
            "frame_color": WARM_GOLD,
            "bullets": [
                "Compaction Compiler: Compacts physical cell boundaries.",
                "Gdstk Exporter: High-performance binary GDSII/OASIS.",
                "WATCHDOG cc_watcher: TCL watchdog script parses JSON exchange.",
                "Native EDA: Imports layouts in Virtuoso/Custom Compiler."
            ]
        }
    ]
    
    for p_idx, p_data in enumerate(panels_data):
        left = Inches(panel_lefts[p_idx])
        top = Inches(1.4)
        width = Inches(panel_width)
        height = Inches(panel_height)
        
        # draw panel
        add_shadow_card(slide, left, top, width, height, MSO_SHAPE.ROUNDED_RECTANGLE)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE_BG
        box.line.color.rgb = p_data["frame_color"]
        box.line.width = Pt(2)
        
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.18)
        tf.margin_top = Inches(0.15)
        tf.margin_right = Inches(0.18)
        
        p_title = tf.paragraphs[0]
        p_title.text = p_data["title"]
        p_title.font.name = "Calibri"
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = p_data["frame_color"]
        p_title.alignment = PP_ALIGN.CENTER
        p_title.space_after = Pt(12)
        
        for b in p_data["bullets"]:
            p = tf.add_paragraph()
            p.text = f"• {b}"
            p.font.name = "Calibri"
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_CHARCOAL
            p.space_after = Pt(5)
            p.level = 0

def add_visual_pipeline(slide, idx, total_pages):
    """Draws the system block diagram flowchart natively using shapes (replaces image)."""
    add_slide_header(slide, "System Architecture: 4-Stage Compilation Block Diagram")
    add_slide_footer(slide, idx, total_pages)
    
    # 1. SPICE Ingest box
    add_shadow_card(slide, Inches(0.6), Inches(3.2), Inches(1.8), Inches(1.3), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_spice = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(3.2), Inches(1.8), Inches(1.3))
    box_spice.fill.solid()
    box_spice.fill.fore_color.rgb = SOFT_CREAM
    box_spice.line.color.rgb = WARM_GOLD
    box_spice.line.width = Pt(1.5)
    tf_sp = box_spice.text_frame
    tf_sp.word_wrap = True
    tf_sp.margin_top = Inches(0.1)
    p_sp = tf_sp.paragraphs[0]
    p_sp.text = "SPICE Netlist\nInput (.sp / .cdl)"
    p_sp.font.name = "Calibri"
    p_sp.font.size = Pt(12)
    p_sp.font.bold = True
    p_sp.font.color.rgb = ACADEMIC_BLUE
    p_sp.alignment = PP_ALIGN.CENTER
    
    # Stage 1 Box
    add_shadow_card(slide, Inches(3.0), Inches(3.2), Inches(2.1), Inches(1.3), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_s1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(3.2), Inches(2.1), Inches(1.3))
    box_s1.fill.solid()
    box_s1.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_s1.line.color.rgb = ACADEMIC_BLUE
    box_s1.line.width = Pt(1.5)
    tf_s1 = box_s1.text_frame
    tf_s1.word_wrap = True
    tf_s1.margin_top = Inches(0.1)
    p_s1 = tf_s1.paragraphs[0]
    p_s1.text = "Stage 1: Parser\n- SPICE Ingestion\n- Subgraph Isomorphism\n- Symmetry Extraction"
    p_s1.font.name = "Calibri"
    p_s1.font.size = Pt(10.5)
    p_s1.font.bold = True
    p_s1.font.color.rgb = ACADEMIC_BLUE
    
    # Stage 2 Box
    add_shadow_card(slide, Inches(5.7), Inches(3.2), Inches(2.1), Inches(1.3), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_s2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.7), Inches(3.2), Inches(2.1), Inches(1.3))
    box_s2.fill.solid()
    box_s2.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_s2.line.color.rgb = ROYAL_BLUE
    box_s2.line.width = Pt(1.5)
    tf_s2 = box_s2.text_frame
    tf_s2.word_wrap = True
    tf_s2.margin_top = Inches(0.1)
    p_s2 = tf_s2.paragraphs[0]
    p_s2.text = "Stage 2: AI Placer\n- LangGraph Planning\n- ReAct DRC Healing\n- Diffusion Compaction"
    p_s2.font.name = "Calibri"
    p_s2.font.size = Pt(10.5)
    p_s2.font.bold = True
    p_s2.font.color.rgb = ROYAL_BLUE
    
    # Stage 4 Box
    add_shadow_card(slide, Inches(8.4), Inches(3.2), Inches(2.1), Inches(1.3), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_s4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.4), Inches(3.2), Inches(2.1), Inches(1.3))
    box_s4.fill.solid()
    box_s4.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_s4.line.color.rgb = CODE_BLUE
    box_s4.line.width = Pt(1.5)
    tf_s4 = box_s4.text_frame
    tf_s4.word_wrap = True
    tf_s4.margin_top = Inches(0.1)
    p_s4 = tf_s4.paragraphs[0]
    p_s4.text = "Stage 4: Exporters\n- Physical Compaction\n- Gdstk OASIS/GDSII\n- cc_watcher Watchdog"
    p_s4.font.name = "Calibri"
    p_s4.font.size = Pt(10.5)
    p_s4.font.bold = True
    p_s4.font.color.rgb = CODE_BLUE
    
    # Output box
    add_shadow_card(slide, Inches(11.1), Inches(3.2), Inches(1.7), Inches(1.3), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_out = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.1), Inches(3.2), Inches(1.7), Inches(1.3))
    box_out.fill.solid()
    box_out.fill.fore_color.rgb = SOFT_CREAM
    box_out.line.color.rgb = WARM_GOLD
    box_out.line.width = Pt(1.5)
    tf_out = box_out.text_frame
    tf_out.word_wrap = True
    tf_out.margin_top = Inches(0.1)
    p_out = tf_out.paragraphs[0]
    p_out.text = "Outputs\n- GDSII / OASIS\n- OpenAccess DB"
    p_out.font.name = "Calibri"
    p_out.font.size = Pt(12)
    p_out.font.bold = True
    p_out.font.color.rgb = ACADEMIC_BLUE
    p_out.alignment = PP_ALIGN.CENTER
    
    # Stage 3 box (GUI co-pilot) above Stage 2
    add_shadow_card(slide, Inches(5.7), Inches(1.3), Inches(2.1), Inches(1.2), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_s3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.7), Inches(1.3), Inches(2.1), Inches(1.2))
    box_s3.fill.solid()
    box_s3.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_s3.line.color.rgb = CODE_GREEN
    box_s3.line.width = Pt(1.5)
    tf_s3 = box_s3.text_frame
    tf_s3.word_wrap = True
    tf_s3.margin_top = Inches(0.1)
    p_s3 = tf_s3.paragraphs[0]
    p_s3.text = "Stage 3: GUI Co-Pilot\n- PySide6 Canvas\n- Intent Classifiers\n- Conversational Panel"
    p_s3.font.name = "Calibri"
    p_s3.font.size = Pt(10.5)
    p_s3.font.bold = True
    p_s3.font.color.rgb = CODE_GREEN
    
    # Connectors
    add_block_arrow(slide, Inches(2.5), Inches(3.65), Inches(0.45), Inches(0.35))
    add_block_arrow(slide, Inches(5.15), Inches(3.65), Inches(0.5), Inches(0.35))
    add_block_arrow(slide, Inches(7.85), Inches(3.65), Inches(0.5), Inches(0.35))
    add_block_arrow(slide, Inches(10.55), Inches(3.65), Inches(0.55), Inches(0.35))
    
    # Bidirectional vertical arrows between Stage 3 (GUI) and Stage 2 (Placer)
    add_block_arrow(slide, Inches(6.3), Inches(2.6), Inches(0.35), Inches(0.5), orient="D", color=CODE_GREEN)
    add_block_arrow(slide, Inches(7.05), Inches(2.6), Inches(0.35), Inches(0.5), orient="U", color=CODE_GREEN)
    
    # Bottom description label box
    desc_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(12.0), Inches(1.2))
    tf_d = desc_box.text_frame
    tf_d.word_wrap = True
    p_d1 = tf_d.paragraphs[0]
    p_d1.text = "The 4-stage pipeline is completely decoupled: the Strategic AI Planner (Stage 2) recommends relative topologies, which the Snapping and Compaction engines (Stage 4) translate into sub-micron geometries."
    p_d1.font.name = "Calibri"
    p_d1.font.size = Pt(12)
    p_d1.font.color.rgb = TEXT_CHARCOAL
    
    p_d2 = tf_d.add_paragraph()
    p_d2.text = "Stage 3 acts as the Human-in-the-loop bridge, parsing chat override commands into XML parameters and flushing them asynchronously to redraw the canvas without freezing the thread."
    p_d2.font.name = "Calibri"
    p_d2.font.size = Pt(12)
    p_d2.font.color.rgb = TEXT_CHARCOAL
    p_d2.space_before = Pt(4)

def add_visual_parser(slide, idx, total_pages):
    """Draws the Stage 1 ingestion and subgraph matching flow natively."""
    add_slide_header(slide, "Stage 1: Parser Ingestion Flowchart")
    add_slide_footer(slide, idx, total_pages)
    
    # Row 1: 3 cards
    lefts_r1 = [0.8, 5.0, 9.2]
    width_box = 3.3
    height_box = 1.6
    
    r1_data = [
        {"title": "1. CDL/SPICE netlist", "color": ROYAL_BLUE, "desc": "Ingests hierarchical netlist files and flattens device sub-module instances."},
        {"title": "2. Properties Resolver", "color": CODE_BLUE, "desc": "Extracts transistor parameters: widths (W), finger multipliers (M), and lengths (L)."},
        {"title": "3. Bipartite Graph", "color": CODE_GREEN, "desc": "Converts netlist to a bipartite graph where nodes are devices/nets and edges are pins."}
    ]
    
    for c_idx, rd in enumerate(r1_data):
        add_shadow_card(slide, Inches(lefts_r1[c_idx]), Inches(1.75), Inches(width_box), Inches(height_box), MSO_SHAPE.ROUNDED_RECTANGLE)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lefts_r1[c_idx]), Inches(1.75), Inches(width_box), Inches(height_box))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE_BG
        box.line.color.rgb = rd["color"]
        box.line.width = Pt(1.5)
        
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.1)
        
        p_title = tf.paragraphs[0]
        p_title.text = rd["title"]
        p_title.font.name = "Calibri"
        p_title.font.size = Pt(14)
        p_title.font.bold = True
        p_title.font.color.rgb = rd["color"]
        
        p_desc = tf.add_paragraph()
        p_desc.text = rd["desc"]
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = TEXT_CHARCOAL
        p_desc.space_before = Pt(3)
        
    # Row 2: 3 cards
    r2_data = [
        {"title": "6. PDK Snap mapping", "color": WARM_GOLD, "desc": "Reads PDK track settings and maps snapping constraints to poly poly gates and horizontal fins."},
        {"title": "5. Constraints Database", "color": RGBColor(230, 90, 40), "desc": "Groups identified devices into hierarchical matching constraint blocks for the Placer state."},
        {"title": "4. Symmetry Detector", "color": CODE_GREEN, "desc": "Executes graph isomorphism walking (NetworkX) to identify matched pairs & current mirrors."}
    ]
    
    # R2 flows backwards: from right to left (c=2 to c=0)
    lefts_r2 = [0.8, 5.0, 9.2]
    
    for c_idx, rd in enumerate(r2_data):
        # We index rd so c_idx=0 is index 2, c_idx=1 is index 1, c_idx=2 is index 0 to flow backwards
        add_shadow_card(slide, Inches(lefts_r2[c_idx]), Inches(4.55), Inches(width_box), Inches(height_box), MSO_SHAPE.ROUNDED_RECTANGLE)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lefts_r2[c_idx]), Inches(4.55), Inches(width_box), Inches(height_box))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE_BG
        box.line.color.rgb = rd["color"]
        box.line.width = Pt(1.5)
        
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.1)
        
        p_title = tf.paragraphs[0]
        p_title.text = rd["title"]
        p_title.font.name = "Calibri"
        p_title.font.size = Pt(14)
        p_title.font.bold = True
        p_title.font.color.rgb = rd["color"]
        
        p_desc = tf.add_paragraph()
        p_desc.text = rd["desc"]
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(10.5)
        p_desc.font.color.rgb = TEXT_CHARCOAL
        p_desc.space_before = Pt(3)
        
    # Arrows connecting cards
    add_block_arrow(slide, Inches(4.2), Inches(2.375), Inches(0.75), Inches(0.35), color=CODE_BLUE)
    add_block_arrow(slide, Inches(8.4), Inches(2.375), Inches(0.75), Inches(0.35), color=CODE_GREEN)
    
    # Downward arrow on right side from Box 3 (top right) to Box 4 (bottom right)
    add_block_arrow(slide, Inches(10.6), Inches(3.45), Inches(0.4), Inches(0.9), orient="D", color=CODE_GREEN)
    
    # Backward arrows in Row 2
    add_block_arrow(slide, Inches(8.4), Inches(5.175), Inches(0.75), Inches(0.35), orient="L", color=RGBColor(230, 90, 40))
    add_block_arrow(slide, Inches(4.2), Inches(5.175), Inches(0.75), Inches(0.35), orient="L", color=WARM_GOLD)

def add_visual_placement_flow(slide, idx, total_pages):
    """Draws a complete native flowchart of the Planner-Critic self-healing loop with decision diamonds."""
    add_slide_header(slide, "Stage 2: Placer-Critic Self-Healing Flowchart")
    add_slide_footer(slide, idx, total_pages)
    
    # 1. Builder Start box
    add_shadow_card(slide, Inches(0.566), Inches(4.25), Inches(2.0), Inches(1.1), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_start = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.566), Inches(4.25), Inches(2.0), Inches(1.1))
    box_start.fill.solid()
    box_start.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_start.line.color.rgb = ROYAL_BLUE
    box_start.line.width = Pt(1.5)
    tf_st = box_start.text_frame
    tf_st.word_wrap = True
    p_st = tf_st.paragraphs[0]
    p_st.text = "Builder Node\nInitializes row & cell structures"
    p_st.font.name = "Calibri"
    p_st.font.size = Pt(10.5)
    p_st.font.bold = True
    p_st.font.color.rgb = ROYAL_BLUE
    p_st.alignment = PP_ALIGN.CENTER
    
    # 2. Specialist ordering box
    add_shadow_card(slide, Inches(3.166), Inches(4.25), Inches(2.0), Inches(1.1), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_spec = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.166), Inches(4.25), Inches(2.0), Inches(1.1))
    box_spec.fill.solid()
    box_spec.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_spec.line.color.rgb = CODE_BLUE
    box_spec.line.width = Pt(1.5)
    tf_sp = box_spec.text_frame
    tf_sp.word_wrap = True
    p_sp = tf_sp.paragraphs[0]
    p_sp.text = "Specialist Node\nProposes relative ordering & splits"
    p_sp.font.name = "Calibri"
    p_sp.font.size = Pt(10.5)
    p_sp.font.bold = True
    p_sp.font.color.rgb = CODE_BLUE
    p_sp.alignment = PP_ALIGN.CENTER
    
    # 3. Spacing Critic check box
    add_shadow_card(slide, Inches(5.766), Inches(4.25), Inches(2.0), Inches(1.1), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_critic = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.766), Inches(4.25), Inches(2.0), Inches(1.1))
    box_critic.fill.solid()
    box_critic.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_critic.line.color.rgb = CODE_GREEN
    box_critic.line.width = Pt(1.5)
    tf_cr = box_critic.text_frame
    tf_cr.word_wrap = True
    p_cr = tf_cr.paragraphs[0]
    p_cr.text = "DRC Critic Node\nSweepline spacing & overlap checks"
    p_cr.font.name = "Calibri"
    p_cr.font.size = Pt(10.5)
    p_cr.font.bold = True
    p_cr.font.color.rgb = CODE_GREEN
    p_cr.alignment = PP_ALIGN.CENTER
    
    # 4. Decision Diamond: Spacing violations?
    add_shadow_card(slide, Inches(8.366), Inches(3.90), Inches(1.8), Inches(1.8), MSO_SHAPE.DIAMOND)
    diamond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(8.366), Inches(3.90), Inches(1.8), Inches(1.8))
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = SOFT_CREAM
    diamond.line.color.rgb = WARM_GOLD
    diamond.line.width = Pt(1.5)
    tf_dia = diamond.text_frame
    tf_dia.word_wrap = True
    p_dia = tf_dia.paragraphs[0]
    p_dia.text = "DRC Spacing\nViolations?"
    p_dia.font.name = "Calibri"
    p_dia.font.size = Pt(10.5)
    p_dia.font.bold = True
    p_dia.font.color.rgb = ACADEMIC_BLUE
    p_dia.alignment = PP_ALIGN.CENTER
    
    # Yes branch: offset calculation (above diamond)
    add_shadow_card(slide, Inches(8.266), Inches(2.25), Inches(2.0), Inches(1.1), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_yes = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.266), Inches(2.25), Inches(2.0), Inches(1.1))
    box_yes.fill.solid()
    box_yes.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_yes.line.color.rgb = RGBColor(230, 90, 40)
    box_yes.line.width = Pt(1.5)
    tf_yes = box_yes.text_frame
    tf_yes.word_wrap = True
    p_yes = tf_yes.paragraphs[0]
    p_yes.text = "Compute coordinate shift offset vectors"
    p_yes.font.name = "Calibri"
    p_yes.font.size = Pt(10.5)
    p_yes.font.bold = True
    p_yes.font.color.rgb = RGBColor(230, 90, 40)
    p_yes.alignment = PP_ALIGN.CENTER
    
    # No branch: physical compaction exporter
    add_shadow_card(slide, Inches(10.766), Inches(4.25), Inches(2.0), Inches(1.1), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_no = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.766), Inches(4.25), Inches(2.0), Inches(1.1))
    box_no.fill.solid()
    box_no.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_no.line.color.rgb = CODE_BLUE
    box_no.line.width = Pt(1.5)
    tf_no = box_no.text_frame
    tf_no.word_wrap = True
    p_no = tf_no.paragraphs[0]
    p_no.text = "Physical Compaction\nProceeds to exporters"
    p_no.font.name = "Calibri"
    p_no.font.size = Pt(10.5)
    p_no.font.bold = True
    p_no.font.color.rgb = CODE_BLUE
    p_no.alignment = PP_ALIGN.CENTER
    
    # Connecting Horizontal Arrows (Centered exactly at Y = 4.625)
    add_block_arrow(slide, Inches(2.666), Inches(4.625), Inches(0.4), Inches(0.35), color=CODE_BLUE)
    add_block_arrow(slide, Inches(5.266), Inches(4.625), Inches(0.4), Inches(0.35), color=CODE_GREEN)
    add_block_arrow(slide, Inches(7.866), Inches(4.625), Inches(0.4), Inches(0.35), color=WARM_GOLD)
    add_block_arrow(slide, Inches(10.266), Inches(4.625), Inches(0.4), Inches(0.35), color=CODE_BLUE)
    
    # Diamond yes upward arrow (Centered exactly at X = 9.266)
    add_block_arrow(slide, Inches(9.091), Inches(3.425), Inches(0.35), Inches(0.4), orient="U", color=RGBColor(230, 90, 40))
    
    # Label "Yes" next to upward arrow
    tx_yes = slide.shapes.add_textbox(Inches(9.4), Inches(3.475), Inches(0.6), Inches(0.3))
    tx_yes.text_frame.paragraphs[0].text = "Yes"
    tx_yes.text_frame.paragraphs[0].font.size = Pt(11)
    tx_yes.text_frame.paragraphs[0].font.bold = True
    tx_yes.text_frame.paragraphs[0].font.color.rgb = RGBColor(230, 90, 40)
    
    # Clean continuous line connector loop back to Specialist:
    # 1. Horizontal segment from left of yes box (X=8.266) to center of Specialist (X=4.166)
    rect_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.166), Inches(2.78), Inches(4.1), Inches(0.04))
    rect_h.fill.solid()
    rect_h.fill.fore_color.rgb = RGBColor(230, 90, 40)
    rect_h.line.fill.background()
    
    # 2. Vertical segment from X=4.166 down to Specialist top Y=4.25
    rect_v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.146), Inches(2.82), Inches(0.04), Inches(1.18))
    rect_v.fill.solid()
    rect_v.fill.fore_color.rgb = RGBColor(230, 90, 40)
    rect_v.line.fill.background()
    
    # 3. Downward arrowhead pointing directly at Specialist box top (starts at Y=4.25)
    arrow_down = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(4.016), Inches(4.00), Inches(0.3), Inches(0.25))
    arrow_down.rotation = 180
    arrow_down.fill.solid()
    arrow_down.fill.fore_color.rgb = RGBColor(230, 90, 40)
    arrow_down.line.fill.background()
    
    # Label "No" above the compaction arrow
    tx_no = slide.shapes.add_textbox(Inches(10.2), Inches(4.2), Inches(0.5), Inches(0.3))
    tx_no.text_frame.paragraphs[0].text = "No"
    tx_no.text_frame.paragraphs[0].font.size = Pt(11)
    tx_no.text_frame.paragraphs[0].font.bold = True
    tx_no.text_frame.paragraphs[0].font.color.rgb = CODE_BLUE


def add_visual_chatbot_flow(slide, idx, total_pages):
    """Draws the Stage 3 chatbot GUI and QThread worker flowchart natively."""
    add_slide_header(slide, "Stage 3: ChatPanel Co-Pilot Event Queue Flowchart")
    add_slide_footer(slide, idx, total_pages)
    
    box_w = 2.6
    box_h = 1.3
    top_r1 = 2.0
    top_r2 = 4.6
    lefts = [0.8, 3.8, 6.8, 9.8]
    
    r1_data = [
        {"title": "1. User Prompt Text", "color": ROYAL_BLUE, "desc": "Designer types layout query text (e.g. \"align gates\")."},
        {"title": "2. PySide6 Canvas GUI", "color": CODE_GREEN, "desc": "ChatPanel captures text; triggers orchestrator workers."},
        {"title": "3. OrchestratorWorker", "color": CODE_BLUE, "desc": "Background worker QThread runs model calculations."},
        {"title": "4. Intent Classifier", "color": WARM_GOLD, "desc": "Classifies prompt into move, swap, align, abut action intents."}
    ]
    
    for c_idx, rd in enumerate(r1_data):
        add_shadow_card(slide, Inches(lefts[c_idx]), Inches(top_r1), Inches(box_w), Inches(box_h), MSO_SHAPE.ROUNDED_RECTANGLE)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lefts[c_idx]), Inches(top_r1), Inches(box_w), Inches(box_h))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE_BG
        box.line.color.rgb = rd["color"]
        box.line.width = Pt(1.5)
        
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_top = Inches(0.08)
        
        p_title = tf.paragraphs[0]
        p_title.text = rd["title"]
        p_title.font.name = "Calibri"
        p_title.font.size = Pt(13)
        p_title.font.bold = True
        p_title.font.color.rgb = rd["color"]
        
        p_desc = tf.add_paragraph()
        p_desc.text = rd["desc"]
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(9.5)
        p_desc.font.color.rgb = TEXT_CHARCOAL
        p_desc.space_before = Pt(2)
        
    r2_data = [
        {"title": "8. Undo Stack Log", "color": ACADEMIC_BLUE, "desc": "Logs layout modifications on the undo stack for reverts."},
        {"title": "7. QGraphicsScene", "color": CODE_BLUE, "desc": "Redraws layout shapes instantly with zero lagging."},
        {"title": "6. QTimer Batch Flusher", "color": CODE_GREEN, "desc": "Groups commands into a single canvas redraw transaction."},
        {"title": "5. XML Command Queue", "color": RGBColor(230, 90, 40), "desc": "Pushes parsed commands to structured queue buffer."}
    ]
    
    for c_idx, rd in enumerate(r2_data):
        add_shadow_card(slide, Inches(lefts[c_idx]), Inches(top_r2), Inches(box_w), Inches(box_h), MSO_SHAPE.ROUNDED_RECTANGLE)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lefts[c_idx]), Inches(top_r2), Inches(box_w), Inches(box_h))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE_BG
        box.line.color.rgb = rd["color"]
        box.line.width = Pt(1.5)
        
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_top = Inches(0.08)
        
        p_title = tf.paragraphs[0]
        p_title.text = rd["title"]
        p_title.font.name = "Calibri"
        p_title.font.size = Pt(13)
        p_title.font.bold = True
        p_title.font.color.rgb = rd["color"]
        
        p_desc = tf.add_paragraph()
        p_desc.text = rd["desc"]
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(9.5)
        p_desc.font.color.rgb = TEXT_CHARCOAL
        p_desc.space_before = Pt(2)
        
    # Arrows Row 1
    add_block_arrow(slide, Inches(3.45), Inches(top_r1 + 0.5), Inches(0.3), Inches(0.25), color=CODE_GREEN)
    add_block_arrow(slide, Inches(6.45), Inches(top_r1 + 0.5), Inches(0.3), Inches(0.25), color=CODE_BLUE)
    add_block_arrow(slide, Inches(9.45), Inches(top_r1 + 0.5), Inches(0.3), Inches(0.25), color=WARM_GOLD)
    
    # Downward arrow to Row 2
    add_block_arrow(slide, Inches(10.9), Inches(3.45), Inches(0.3), Inches(1.05), orient="D", color=RGBColor(230, 90, 40))
    
    # Arrows Row 2 (backwards: from right to left)
    add_block_arrow(slide, Inches(9.45), Inches(top_r2 + 0.5), Inches(0.3), Inches(0.25), orient="L", color=CODE_GREEN)
    add_block_arrow(slide, Inches(6.45), Inches(top_r2 + 0.5), Inches(0.3), Inches(0.25), orient="L", color=CODE_BLUE)
    add_block_arrow(slide, Inches(3.45), Inches(top_r2 + 0.5), Inches(0.3), Inches(0.25), orient="L", color=ACADEMIC_BLUE)

def add_visual_exporter_flow(slide, idx, total_pages):
    """Draws the Stage 4 spatial compaction and dual exporter paths flowchart natively."""
    add_slide_header(slide, "Stage 4: Compaction & Dual Exporter Flowchart")
    add_slide_footer(slide, idx, total_pages)
    
    # 1. Start Box: Snapped coordinates
    add_shadow_card(slide, Inches(0.8), Inches(3.4), Inches(2.2), Inches(1.3), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_start = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.4), Inches(2.2), Inches(1.3))
    box_start.fill.solid()
    box_start.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_start.line.color.rgb = ROYAL_BLUE
    box_start.line.width = Pt(1.5)
    tf_st = box_start.text_frame
    tf_st.word_wrap = True
    tf_st.margin_top = Inches(0.1)
    p_st = tf_st.paragraphs[0]
    p_st.text = "Snapped Coordinates\n- Transistor rows\n- Gate poly centers\n- Shared diffusions"
    p_st.font.name = "Calibri"
    p_st.font.size = Pt(10.5)
    p_st.font.bold = True
    p_st.font.color.rgb = ROYAL_BLUE
    
    # 2. Compaction Engine Box
    add_shadow_card(slide, Inches(3.7), Inches(3.4), Inches(2.2), Inches(1.3), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_comp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.7), Inches(3.4), Inches(2.2), Inches(1.3))
    box_comp.fill.solid()
    box_comp.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_comp.line.color.rgb = ACADEMIC_BLUE
    box_comp.line.width = Pt(1.5)
    tf_cp = box_comp.text_frame
    tf_cp.word_wrap = True
    tf_cp.margin_top = Inches(0.1)
    p_cp = tf_cp.paragraphs[0]
    p_cp.text = "Compaction Engine\n- Runs sweep clearances\n- Optimizes cell bounds\n- Aligns routing pins"
    p_cp.font.name = "Calibri"
    p_cp.font.size = Pt(10.5)
    p_cp.font.bold = True
    p_cp.font.color.rgb = ACADEMIC_BLUE
    
    # Branch 1 (top): Gdstk Exporter
    add_shadow_card(slide, Inches(6.8), Inches(1.9), Inches(2.2), Inches(1.2), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_gdstk = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.9), Inches(2.2), Inches(1.2))
    box_gdstk.fill.solid()
    box_gdstk.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_gdstk.line.color.rgb = CODE_BLUE
    box_gdstk.line.width = Pt(1.5)
    tf_gd = box_gdstk.text_frame
    tf_gd.word_wrap = True
    tf_gd.margin_top = Inches(0.1)
    p_gd = tf_gd.paragraphs[0]
    p_gd.text = "Gdstk Exporter\nClones PDK base cells;\ncompiles geometries"
    p_gd.font.name = "Calibri"
    p_gd.font.size = Pt(10.5)
    p_gd.font.bold = True
    p_gd.font.color.rgb = CODE_BLUE
    
    # Branch 1 end: Binary GDSII/OASIS
    add_shadow_card(slide, Inches(9.8), Inches(1.9), Inches(2.6), Inches(1.2), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_bin = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(1.9), Inches(2.6), Inches(1.2))
    box_bin.fill.solid()
    box_bin.fill.fore_color.rgb = SOFT_CREAM
    box_bin.line.color.rgb = WARM_GOLD
    box_bin.line.width = Pt(1.5)
    tf_bi = box_bin.text_frame
    tf_bi.word_wrap = True
    tf_bi.margin_top = Inches(0.1)
    p_bi = tf_bi.paragraphs[0]
    p_bi.text = "Binary GDSII / OASIS Streams\nHighly compacted file size;\nready for verification/tape-out"
    p_bi.font.name = "Calibri"
    p_bi.font.size = Pt(10.5)
    p_bi.font.bold = True
    p_bi.font.color.rgb = ACADEMIC_BLUE
    
    # Branch 2 (bottom): JSON exchange file
    add_shadow_card(slide, Inches(6.8), Inches(4.8), Inches(2.2), Inches(1.2), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_json = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.8), Inches(2.2), Inches(1.2))
    box_json.fill.solid()
    box_json.fill.fore_color.rgb = SOFT_CREAM
    box_json.line.color.rgb = WARM_GOLD
    box_json.line.width = Pt(1.5)
    tf_js = box_json.text_frame
    tf_js.word_wrap = True
    tf_js.margin_top = Inches(0.1)
    p_js = tf_js.paragraphs[0]
    p_js.text = "layout.json file\nStores snapped coords;\nlocalized exchange buffer"
    p_js.font.name = "Calibri"
    p_js.font.size = Pt(10.5)
    p_js.font.bold = True
    p_js.font.color.rgb = ACADEMIC_BLUE
    
    # Branch 2 end: cc_watcher.tcl Watchdog
    add_shadow_card(slide, Inches(9.8), Inches(4.8), Inches(2.6), Inches(1.2), MSO_SHAPE.ROUNDED_RECTANGLE)
    box_watch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(4.8), Inches(2.6), Inches(1.2))
    box_watch.fill.solid()
    box_watch.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_watch.line.color.rgb = RGBColor(230, 90, 40)
    box_watch.line.width = Pt(1.5)
    tf_wa = box_watch.text_frame
    tf_wa.word_wrap = True
    tf_wa.margin_top = Inches(0.1)
    p_wa = tf_wa.paragraphs[0]
    p_wa.text = "cc_watcher Watchdog\nInjects coordinate changes;\nupdates Virtuoso cellview"
    p_wa.font.name = "Calibri"
    p_wa.font.size = Pt(10.5)
    p_wa.font.bold = True
    p_wa.font.color.rgb = RGBColor(230, 90, 40)
    
    # Arrows
    add_block_arrow(slide, Inches(3.1), Inches(3.875), Inches(0.5), Inches(0.35), color=ACADEMIC_BLUE)
    
    # Split arrows
    # To Gdstk (top branch)
    arrow_top = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.0), Inches(3.7), Inches(0.6), Inches(0.35))
    arrow_top.fill.solid()
    arrow_top.fill.fore_color.rgb = CODE_BLUE
    arrow_top.line.fill.background()
    # Rotate arrow upward slightly
    arrow_top.rotation = -25
    
    # To JSON (bottom branch)
    arrow_bot = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.0), Inches(4.2), Inches(0.6), Inches(0.35))
    arrow_bot.fill.solid()
    arrow_bot.fill.fore_color.rgb = WARM_GOLD
    arrow_bot.line.fill.background()
    arrow_bot.rotation = 25
    
    # Arrow to Binary
    add_block_arrow(slide, Inches(9.1), Inches(2.325), Inches(0.6), Inches(0.35), color=CODE_BLUE)
    # Arrow to Watchdog
    add_block_arrow(slide, Inches(9.1), Inches(5.225), Inches(0.6), Inches(0.35), color=RGBColor(230, 90, 40))

def add_visual_langgraph(slide, idx, total_pages):
    """Draws a complete native flowchart of the LangGraph Multi-Agent AI Placement Pipeline."""
    add_slide_header(slide, "Stage 2: LangGraph Multi-Agent Planning Pipeline")
    add_slide_footer(slide, idx, total_pages)
    
    # Row 1: 5 cards
    box_w = Inches(2.0)
    box_h = Inches(1.0)
    top_r1 = Inches(1.8)
    top_r2 = Inches(3.9)
    
    # 1. Input Node JSON
    add_shadow_card(slide, Inches(0.6), top_r1, Inches(1.8), box_h, MSO_SHAPE.ROUNDED_RECTANGLE)
    box_start = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), top_r1, Inches(1.8), box_h)
    box_start.fill.solid()
    box_start.fill.fore_color.rgb = SOFT_CREAM
    box_start.line.color.rgb = WARM_GOLD
    box_start.line.width = Pt(1.5)
    tf = box_start.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Input JSON\nLogical device netlist & PDK metrics"
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ACADEMIC_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # 2. Analyst Agent
    add_shadow_card(slide, Inches(3.0), top_r1, Inches(2.0), box_h, MSO_SHAPE.ROUNDED_RECTANGLE)
    box_analyst = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), top_r1, Inches(2.0), box_h)
    box_analyst.fill.solid()
    box_analyst.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_analyst.line.color.rgb = ROYAL_BLUE
    box_analyst.line.width = Pt(1.5)
    tf = box_analyst.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Topology Analyst\nBipartite graph matching of diff-pairs & mirrors"
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # 3. Selector Agent
    add_shadow_card(slide, Inches(5.6), top_r1, Inches(2.0), box_h, MSO_SHAPE.ROUNDED_RECTANGLE)
    box_selector = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.6), top_r1, Inches(2.0), box_h)
    box_selector.fill.solid()
    box_selector.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_selector.line.color.rgb = CODE_BLUE
    box_selector.line.width = Pt(1.5)
    tf = box_selector.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Strategy Selector\nPlans rows, layout symmetry, and centroids"
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CODE_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # 4. Specialist Agent
    add_shadow_card(slide, Inches(8.2), top_r1, Inches(2.1), box_h, MSO_SHAPE.ROUNDED_RECTANGLE)
    box_specialist = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), top_r1, Inches(2.1), box_h)
    box_specialist.fill.solid()
    box_specialist.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_specialist.line.color.rgb = ACADEMIC_BLUE
    box_specialist.line.width = Pt(1.5)
    tf = box_specialist.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Placement Specialist\nGenerates strategic relative floorplans"
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ACADEMIC_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # 5. Finger Expander
    add_shadow_card(slide, Inches(10.9), top_r1, Inches(1.8), box_h, MSO_SHAPE.ROUNDED_RECTANGLE)
    box_expander = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.9), top_r1, Inches(1.8), box_h)
    box_expander.fill.solid()
    box_expander.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_expander.line.color.rgb = CODE_GREEN
    box_expander.line.width = Pt(1.5)
    tf = box_expander.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Finger Expander\nUnrolls fingers & calculates S/D abutment"
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CODE_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    # Row 2 (flows backwards)
    # 6. Routing Previewer
    add_shadow_card(slide, Inches(10.9), top_r2, Inches(1.8), box_h, MSO_SHAPE.ROUNDED_RECTANGLE)
    box_routing = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.9), top_r2, Inches(1.8), box_h)
    box_routing.fill.solid()
    box_routing.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_routing.line.color.rgb = WARM_GOLD
    box_routing.line.width = Pt(1.5)
    tf = box_routing.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Routing Previewer\nHPWL cost & crossing minimization"
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WARM_GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # 7. DRC Critic
    add_shadow_card(slide, Inches(8.2), top_r2, Inches(2.1), box_h, MSO_SHAPE.ROUNDED_RECTANGLE)
    box_critic = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), top_r2, Inches(2.1), box_h)
    box_critic.fill.solid()
    box_critic.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_critic.line.color.rgb = RGBColor(230, 90, 40)
    box_critic.line.width = Pt(1.5)
    tf = box_critic.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "DRC Critic Agent\nSweepline check; extracts violations"
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(230, 90, 40)
    p.alignment = PP_ALIGN.CENTER
    
    # 8. Output DRC Clean JSON
    add_shadow_card(slide, Inches(5.6), top_r2, Inches(2.0), box_h, MSO_SHAPE.ROUNDED_RECTANGLE)
    box_output = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.6), top_r2, Inches(2.0), box_h)
    box_output.fill.solid()
    box_output.fill.fore_color.rgb = SOFT_CREAM
    box_output.line.color.rgb = CODE_GREEN
    box_output.line.width = Pt(1.5)
    tf = box_output.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "DRC-Clean Layout\nExports unified node coordinate maps"
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CODE_GREEN
    p.alignment = PP_ALIGN.CENTER
    
    # 9. PySide6 Canvas
    add_shadow_card(slide, Inches(3.0), top_r2, Inches(2.0), box_h, MSO_SHAPE.ROUNDED_RECTANGLE)
    box_canvas = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), top_r2, Inches(2.0), box_h)
    box_canvas.fill.solid()
    box_canvas.fill.fore_color.rgb = LIGHT_BLUE_BG
    box_canvas.line.color.rgb = ROYAL_BLUE
    box_canvas.line.width = Pt(1.5)
    tf = box_canvas.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PySide6 Canvas\nRedraws dynamic layout geometries"
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Connectors Row 1
    add_block_arrow(slide, Inches(2.5), Inches(1.8 + 0.32), Inches(0.4), Inches(0.35), color=ROYAL_BLUE)
    add_block_arrow(slide, Inches(5.1), Inches(1.8 + 0.32), Inches(0.4), Inches(0.35), color=CODE_BLUE)
    add_block_arrow(slide, Inches(7.7), Inches(1.8 + 0.32), Inches(0.4), Inches(0.35), color=ACADEMIC_BLUE)
    add_block_arrow(slide, Inches(10.4), Inches(1.8 + 0.32), Inches(0.4), Inches(0.35), color=CODE_GREEN)
    
    # Downward arrow to Row 2
    add_block_arrow(slide, Inches(11.62), Inches(2.9), Inches(0.35), Inches(0.9), orient="D", color=WARM_GOLD)
    
    # Connectors Row 2 (flows backwards)
    add_block_arrow(slide, Inches(10.4), Inches(3.9 + 0.32), Inches(0.4), Inches(0.35), orient="L", color=WARM_GOLD)
    add_block_arrow(slide, Inches(7.7), Inches(3.9 + 0.32), Inches(0.4), Inches(0.35), orient="L", color=CODE_GREEN)
    add_block_arrow(slide, Inches(5.1), Inches(3.9 + 0.32), Inches(0.4), Inches(0.35), orient="L", color=ROYAL_BLUE)
    
    # Upward Feedback Loop: DRC Critic (7) -> Specialist (4)
    add_block_arrow(slide, Inches(9.07), Inches(2.9), Inches(0.35), Inches(0.9), orient="U", color=RGBColor(230, 90, 40))
    
    # Label "Feedback"
    tx_fb = slide.shapes.add_textbox(Inches(9.45), Inches(3.1), Inches(1.3), Inches(0.5))
    tf_fb = tx_fb.text_frame
    tf_fb.word_wrap = True
    p_fb = tf_fb.paragraphs[0]
    p_fb.text = "Spacing Violations\nFeedback"
    p_fb.font.name = "Calibri"
    p_fb.font.size = Pt(8.5)
    p_fb.font.bold = True
    p_fb.font.color.rgb = RGBColor(230, 90, 40)
    
    # Label "DRC Clean"
    tx_ok = slide.shapes.add_textbox(Inches(7.6), Inches(4.35), Inches(0.6), Inches(0.3))
    tx_ok.text_frame.paragraphs[0].text = "Clean"
    tx_ok.text_frame.paragraphs[0].font.size = Pt(8.5)
    tx_ok.text_frame.paragraphs[0].font.bold = True
    tx_ok.text_frame.paragraphs[0].font.color.rgb = CODE_GREEN
    
    # Bottom description
    desc_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.3), Inches(12.133), Inches(1.6))
    tf_d = desc_box.text_frame
    tf_d.word_wrap = True
    p_d1 = tf_d.paragraphs[0]
    p_d1.text = "• Multi-Agent Orchestration: LangGraph manages layout generation as a state-machine. A shared state dictionary houses the netlist bipartite graph, matched symmetry constraint groups, and current coordinate matrices."
    p_d1.font.name = "Calibri"
    p_d1.font.size = Pt(11.5)
    p_d1.font.color.rgb = TEXT_CHARCOAL
    
    p_d2 = tf_d.add_paragraph()
    p_d2.text = "• Closed-Loop Convergence: When the DRC Critic detects overlaps or spacing failures, it writes a list of bounding box corrective shift directives back to the Specialist agent, achieving zero-DRC convergence in under 4 loops."
    p_d2.font.name = "Calibri"
    p_d2.font.size = Pt(11.5)
    p_d2.font.color.rgb = TEXT_CHARCOAL
    p_d2.space_before = Pt(4)

def add_visual_nodes(slide, idx, total_pages):
    """Draws a beautiful 3x2 grid of cards describing the LangGraph specialist nodes."""
    add_slide_header(slide, "Stage 2: Specialist Agent Roles & Responsibilities")
    add_slide_footer(slide, idx, total_pages)
    
    lefts = [0.8, 4.8, 8.8]
    tops = [1.4, 4.1]
    
    nodes_info = [
        {
            "title": "Topology Analyst",
            "color": ROYAL_BLUE,
            "bullets": [
                "Bipartite Graph Walks: Parses netlist connectivity using NetworkX graphs.",
                "Symmetry Matching: Automatically detects differential pairs, current mirrors, and matched terminals."
            ]
        },
        {
            "title": "Strategy Selector",
            "color": ACADEMIC_BLUE,
            "bullets": [
                "Floorplan Layouts: Allocates active cells to horizontal rows and maps well guard ring boundaries.",
                "Centroid Axis Planning: Selects horizontal and vertical symmetry reflection axes."
            ]
        },
        {
            "title": "Placement Specialist",
            "color": CODE_BLUE,
            "bullets": [
                "LLM Reasoning Planner: Determines strategic relative orderings of transistors without absolute math.",
                "Relative Commands: Issues topological move, swap, align, and abut directives."
            ]
        },
        {
            "title": "Finger Expander",
            "color": CODE_GREEN,
            "bullets": [
                "Transistor Splitter: Unrolls wide gates into parallel fingers to fit PDK row height constraints.",
                "Active Area Abutment: Maximizes shared source/drain contacts, compressing spacing to 70nm."
            ]
        },
        {
            "title": "Routing Previewer",
            "color": WARM_GOLD,
            "bullets": [
                "HPWL Routing Cost: Minimizes interconnect distance by calculating Half-Parameter Wire Lengths.",
                "Same-Layer Crossing Checks: Optimizes cell orders to avoid signal trace intersections."
            ]
        },
        {
            "title": "DRC Critic Agent",
            "color": RGBColor(230, 90, 40),
            "bullets": [
                "Sweepline Spacing Engine: Evaluates polygon overlaps against PDK clearance lookup rules.",
                "Correction Directives: Formulates shift vector requirements (e.g. shift X by 0.3um) to heal state."
            ]
        }
    ]
    
    card_w = Inches(3.75)
    card_h = Inches(2.4)
    
    for i, info in enumerate(nodes_info):
        r = i // 3
        c = i % 3
        
        add_shadow_card(slide, Inches(lefts[c]), Inches(tops[r]), card_w, card_h, MSO_SHAPE.ROUNDED_RECTANGLE)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lefts[c]), Inches(tops[r]), card_w, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE_BG
        box.line.color.rgb = info["color"]
        box.line.width = Pt(1.5)
        
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.1)
        tf.margin_right = Inches(0.15)
        
        p_t = tf.paragraphs[0]
        p_t.text = info["title"]
        p_t.font.name = "Calibri"
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = info["color"]
        p_t.space_after = Pt(8)
        
        for b in info["bullets"]:
            p = tf.add_paragraph()
            p.text = f"• {b}"
            p.font.name = "Calibri"
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_CHARCOAL
            p.space_after = Pt(4)

def add_visual_routing(slide, idx, total_pages):
    """Draws a flowchart describing the Stage 3 co-pilot intent classifier and command queue."""
    add_slide_header(slide, "Stage 3: ChatPanel Co-Pilot Event Queue Flowchart")
    add_slide_footer(slide, idx, total_pages)
    
    # 2 rows of 4 cards
    lefts = [0.8, 3.8, 6.8, 9.8]
    top_r1 = 1.8
    top_r2 = 4.3
    card_w = Inches(2.7)
    card_h = Inches(1.1)
    
    r1_data = [
        {"title": "1. Natural Language Input", "color": ROYAL_BLUE, "desc": "Designer types layout query text (e.g. \"align gates\")."},
        {"title": "2. QThread Worker", "color": CODE_GREEN, "desc": "ChatPanel captures text; triggers background thread worker."},
        {"title": "3. Intent Router LLM", "color": CODE_BLUE, "desc": "Analyzes prompt against available layout skills (move, abut)."},
        {"title": "4. XML CMD Synthesizer", "color": WARM_GOLD, "desc": "Converts natural intent into structured command parameters."}
    ]
    
    for c_idx, rd in enumerate(r1_data):
        add_shadow_card(slide, Inches(lefts[c_idx]), Inches(top_r1), card_w, card_h, MSO_SHAPE.ROUNDED_RECTANGLE)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lefts[c_idx]), Inches(top_r1), card_w, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE_BG
        box.line.color.rgb = rd["color"]
        box.line.width = Pt(1.5)
        
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_top = Inches(0.08)
        
        p_title = tf.paragraphs[0]
        p_title.text = rd["title"]
        p_title.font.name = "Calibri"
        p_title.font.size = Pt(12)
        p_title.font.bold = True
        p_title.font.color.rgb = rd["color"]
        
        p_desc = tf.add_paragraph()
        p_desc.text = rd["desc"]
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(9)
        p_desc.font.color.rgb = TEXT_CHARCOAL
        p_desc.space_before = Pt(2)
        
    r2_data = [
        {"title": "8. KLayout Live Preview", "color": ACADEMIC_BLUE, "desc": "Synchronizes coordinates; displays unrolled shapes instantly."},
        {"title": "7. PySide6 Canvas View", "color": CODE_BLUE, "desc": "Redraws layout shapes on the symbolic canvas view."},
        {"title": "6. QGraphicsScene Stack", "color": CODE_GREEN, "desc": "Translates commands into QGraphicsItems & commits to undo/redo."},
        {"title": "5. XML Execution Queue", "color": RGBColor(230, 90, 40), "desc": "Pushes parsed commands to a structured transaction queue."}
    ]
    
    for c_idx, rd in enumerate(r2_data):
        add_shadow_card(slide, Inches(lefts[c_idx]), Inches(top_r2), card_w, card_h, MSO_SHAPE.ROUNDED_RECTANGLE)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lefts[c_idx]), Inches(top_r2), card_w, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE_BG
        box.line.color.rgb = rd["color"]
        box.line.width = Pt(1.5)
        
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.1)
        tf.margin_top = Inches(0.08)
        
        p_title = tf.paragraphs[0]
        p_title.text = rd["title"]
        p_title.font.name = "Calibri"
        p_title.font.size = Pt(12)
        p_title.font.bold = True
        p_title.font.color.rgb = rd["color"]
        
        p_desc = tf.add_paragraph()
        p_desc.text = rd["desc"]
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(9)
        p_desc.font.color.rgb = TEXT_CHARCOAL
        p_desc.space_before = Pt(2)
        
    # Arrows Row 1
    add_block_arrow(slide, Inches(3.55), Inches(top_r1 + 0.4), Inches(0.22), Inches(0.25), color=CODE_GREEN)
    add_block_arrow(slide, Inches(6.55), Inches(top_r1 + 0.4), Inches(0.22), Inches(0.25), color=CODE_BLUE)
    add_block_arrow(slide, Inches(9.55), Inches(top_r1 + 0.4), Inches(0.22), Inches(0.25), color=WARM_GOLD)
    
    # Downward arrow to Row 2
    add_block_arrow(slide, Inches(10.9), Inches(2.95), Inches(0.25), Inches(1.3), orient="D", color=RGBColor(230, 90, 40))
    
    # Arrows Row 2 (backwards: from right to left)
    add_block_arrow(slide, Inches(9.55), Inches(top_r2 + 0.4), Inches(0.22), Inches(0.25), orient="L", color=CODE_GREEN)
    add_block_arrow(slide, Inches(6.55), Inches(top_r2 + 0.4), Inches(0.22), Inches(0.25), orient="L", color=CODE_BLUE)
    add_block_arrow(slide, Inches(3.55), Inches(top_r2 + 0.4), Inches(0.22), Inches(0.25), orient="L", color=ACADEMIC_BLUE)
    
    # Explanatory Text Box below flowchart
    tx_info = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.3))
    tf_i = tx_info.text_frame
    tf_i.word_wrap = True
    p_i = tf_i.paragraphs[0]
    p_i.text = "• Thread Isolation: By offloading heavy model queries and routing calculations to a separate background QThread, the main PySide6 canvas view remains responsive with zero UI freezing."
    p_i.font.name = "Calibri"
    p_i.font.size = Pt(11)
    p_i.font.color.rgb = TEXT_CHARCOAL
    
    p_i2 = tf_i.add_paragraph()
    p_i2.text = "• Transactional Safety: Commands in the queue are validated against snap pitches and compiled as a batch. A lightweight QTimer flusher, reducing canvas repainting overhead."
    p_i2.font.name = "Calibri"
    p_i2.font.size = Pt(11)
    p_i2.font.color.rgb = TEXT_CHARCOAL
    p_i2.space_before = Pt(4)

def add_visual_watchdog(slide, idx, total_pages):
    """Draws a flowchart describing the Stage 4 OpenAccess Database Watchdog observer."""
    add_slide_header(slide, "Stage 4: OpenAccess DB Watchdog Flowchart")
    add_slide_footer(slide, idx, total_pages)
    
    # 5-step horizontal pipeline
    lefts = [0.8, 3.2, 5.7, 8.2, 10.7]
    card_w = Inches(2.0)
    card_h = Inches(1.5)
    
    steps_data = [
        {"title": "1. Exporter Output", "color": ROYAL_BLUE, "desc": "Compactor writes snapped coordinates to layout.json file on disk."},
        {"title": "2. Filesystem Watchdog", "color": ACADEMIC_BLUE, "desc": "cc_watcher.tcl running inside Virtuoso monitors directory for updates."},
        {"title": "3. TCL Ingestion Parser", "color": CODE_BLUE, "desc": "Reads new layout.json file and decodes device-coordinate pairs."},
        {"title": "4. OpenAccess DB API", "color": CODE_GREEN, "desc": "Directly moves cellview shapes using db_move_device calls in-memory."},
        {"title": "5. Active Canvas Refresh", "color": WARM_GOLD, "desc": "Virtuoso/Custom Compiler layout window redraws coordinates instantly."}
    ]
    
    for c_idx, sd in enumerate(steps_data):
        add_shadow_card(slide, Inches(lefts[c_idx]), Inches(2.5), card_w, card_h, MSO_SHAPE.ROUNDED_RECTANGLE)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lefts[c_idx]), Inches(2.5), card_w, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE_BG
        box.line.color.rgb = sd["color"]
        box.line.width = Pt(1.5)
        
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_top = Inches(0.08)
        tf.margin_right = Inches(0.12)
        
        p_title = tf.paragraphs[0]
        p_title.text = sd["title"]
        p_title.font.name = "Calibri"
        p_title.font.size = Pt(12)
        p_title.font.bold = True
        p_title.font.color.rgb = sd["color"]
        
        p_desc = tf.add_paragraph()
        p_desc.text = sd["desc"]
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(9)
        p_desc.font.color.rgb = TEXT_CHARCOAL
        p_desc.space_before = Pt(3)
        
        # Connectors between horizontal cards
        if c_idx < 4:
            add_block_arrow(slide, Inches(lefts[c_idx] + 2.05), Inches(3.075), Inches(0.3), Inches(0.35), color=sd["color"])
            
    # Description text box below flowchart
    tx_info = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.3))
    tf_i = tx_info.text_frame
    tf_i.word_wrap = True
    p_i = tf_i.paragraphs[0]
    p_i.text = "• Independent Synchronization: Rather than establishing heavy inter-process communication sockets, the system uses a robust file watchdog (cc_watcher.tcl) that polls a directory and reads files based on modification times (mtime)."
    p_i.font.name = "Calibri"
    p_i.font.size = Pt(11.5)
    p_i.font.color.rgb = TEXT_CHARCOAL
    
    p_i2 = tf_i.add_paragraph()
    p_i2.text = "• Direct In-Memory Manipulation: The TCL watchdog runs directly inside the active EDA session (Synopsys Custom Compiler or Cadence Virtuoso). It queries the OpenAccess cellview database, moves transistor instances in-memory, and updates PCell abutment flags programmatically. This ensures DRC-clean layouts load instantly without database locks."
    p_i2.font.name = "Calibri"
    p_i2.font.size = Pt(11.5)
    p_i2.font.color.rgb = TEXT_CHARCOAL
    p_i2.space_before = Pt(5)

def add_future_work_slide(slide, idx, total_pages):
    """Draws a beautiful card-based slide for Section 8: Conclusions & Future Work."""
    add_slide_header(slide, "Conclusions & Future Work: Technical Roadmap")
    add_slide_footer(slide, idx, total_pages)
    
    lefts = [0.8, 4.8, 8.8]
    card_w = Inches(3.7)
    card_h = Inches(5.3)
    
    future_data = [
        {
            "title": "1. AI-Guided Routing Agent",
            "color": ROYAL_BLUE,
            "bullets": [
                "PDK Metal Track Snapping: Develop an AI router specialist that routes net paths on discrete metal layers.",
                "Heuristic Routing: Uses A* search and maze routing guided by LLM-defined wire cost constraints.",
                "Obstacle Clearance: Calculates optimal via placements and avoids pre-existing well boundaries."
            ]
        },
        {
            "title": "2. Thermal-Aware Layouts",
            "color": WARM_GOLD,
            "bullets": [
                "Thermal Mismatch Planning: Models silicon hot-spots and schedules symmetrical devices to cancel mismatch.",
                "2D Centroid Architectures: Proposes 2D common-centroid cross-coupled placement arrays.",
                "Parasitic Balancing: Ensures identical trace lengths to balance active drain capacitances."
            ]
        },
        {
            "title": "3. Nanosheet vertical stacking",
            "color": CODE_GREEN,
            "bullets": [
                "Nanosheet / CFET: Support vertical 3D transistor architectures where NMOS and PMOS gates stack vertically.",
                "Vertical track limits: Extends snapping grid coordinates to include vertical routing clearances.",
                "Automated Sub-5nm: Formulates compaction rules for sub-5nm restricted design rules."
            ]
        }
    ]
    
    for c_idx, fd in enumerate(future_data):
        add_shadow_card(slide, Inches(lefts[c_idx]), Inches(1.4), card_w, card_h, MSO_SHAPE.ROUNDED_RECTANGLE)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lefts[c_idx]), Inches(1.4), card_w, card_h)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE_BG
        box.line.color.rgb = fd["color"]
        box.line.width = Pt(2)
        
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.18)
        tf.margin_top = Inches(0.15)
        tf.margin_right = Inches(0.18)
        
        p_title = tf.paragraphs[0]
        p_title.text = fd["title"]
        p_title.font.name = "Calibri"
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = fd["color"]
        p_title.alignment = PP_ALIGN.CENTER
        p_title.space_after = Pt(12)
        
        for b in fd["bullets"]:
            p = tf.add_paragraph()
            p.text = f"• {b}"
            p.font.name = "Calibri"
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_CHARCOAL
            p.space_after = Pt(6)

def create_presentation():

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # ---------------------------------------------------------
    # DEFINE SLIDE CONTENT DATA
    # ---------------------------------------------------------
    slides_data = [
        # Slide 1: Cover Slide
        {
            "type": "cover",
            "title": "AI-Based Analog Layout Automation: A Conversational, Multi-Agent Framework",
            "logo_aus": "ASU.png",
            "logo_si": "SI.png"
        },
        # Slide 2: Outline (rendered as visual timeline)
        {
            "type": "timeline"
        },
        # Slide 3: Section 1 Title
        {
            "type": "section",
            "title": "Section 1: Introduction & Motivation",
            "subtitle": "The Layout Bottleneck, PDK Complexity, and Automation Challenges"
        },
        # Slide 4: Revolution slide (the complexity-to-symbolic abstraction workflow)
        {
            "type": "revolution"
        },
        # Slide 5: Intro - Manual Layout Bottleneck
        {
            "type": "split",
            "title": "Introduction: Manual Analog Layout Bottleneck",
            "box1_title": "The Traditional Analog Layout Flow",
            "box1_bullets": [
                "Human-Dominated Process: Translation of schematic netlist to layout geometries relies heavily on manual expert drawing.",
                "High Complexity: Requires strict compliance with design rules (DRC), device symmetry matching, and parasitic minimization.",
                "Time-Consuming: Layout phase consumes up to 50% of the entire IC design cycle, slowing time-to-market."
            ],
            "box2_title": "Need for Automation & Symmetry",
            "box2_bullets": [
                "Cycle Compression: Shifting to automated layout can compress design cycles from weeks to minutes.",
                "Matching Crucial: Symmetrical placements (e.g. diff pairs) are mandatory to cancel spatial process variations.",
                "Interactive Interface: Engineers must be able to specify spatial constraints without writing code."
            ]
        },
        # Slide 6: Intro - Challenges in Modern PDKs
        {
            "type": "split",
            "title": "Introduction: Challenges in Modern PDKs",
            "box1_title": "Restricted Design Rules (RDRs)",
            "box1_bullets": [
                "Sub-20nm PDKs: Deep sub-micron nodes impose strict spacing limits on active diffusions, poly gates, and metals.",
                "Guard Rings: Boundary guard rings are mandatory to avoid latch-up and shield active transistors from bulk noise.",
                "Snapping Constraints: Poly gates and fins must align to discrete pitches, restricting spacing configurations."
            ],
            "box2_title": "Layout-Dependent Effects (LDE)",
            "box2_bullets": [
                "Well Proximity Effect (WPE): Ion scatter off well edges changes transistor threshold voltage depending on well distance.",
                "Shallow Trench Isolation (STI): Compressive mechanical stress on active active-area alters channel mobilities based on isolation edge proximity.",
                "Impact: Compaction must minimize LDE while maintaining high packing density."
            ]
        },
        # Slide 7: Intro - Limitations of Existing Automation
        {
            "type": "table_only",
            "title": "Introduction: Limitations of Existing Approaches",
            "headers": ["Approach", "Core Methodology", "Primary Limitations"],
            "rows": [
                ["Procedural Generators", "Hard-coded scripts for specific cells (e.g., Berkeley Analog Generator)", "High development and maintenance overhead; zero portability across PDK nodes."],
                ["Analytical Placers", "Analytical cost function optimization (Simulated Annealing, ILP, Force-Directed)", "Black-box execution; fails to incorporate qualitative designer intent and preferences."],
                ["GenAI Layout Models", "Textual layout tokenizers (e.g., GLayout, LLM-coordinate output)", "Highly susceptible to coordinate arithmetic hallucinations and PDK spacing rounding errors."]
            ]
        },
        # Slide 8: Intro - Proposed Solution
        {
            "type": "split",
            "title": "Introduction: The Proposed Solution",
            "box1_title": "Decoupled Strategic & Geometric Architecture",
            "box1_bullets": [
                "Strategic Planning: Powered by LLM agents. Reasons about symbolic relative ordering (ordering, orientations, finger splits) rather than raw dimensions.",
                "Deterministic Compiler: Snaps symbolic placements to physical sub-micron tracks, enforcing PDK spacing and boundary checks.",
                "Human-in-the-Loop: A conversational PySide6 editor panel dynamically captures designer intent, bridging AI and canvas."
            ],
            "box2_title": "Closed-Loop Correctness",
            "box2_bullets": [
                "Sweepline Spacing Critic: A geometry engine evaluates relative placements against PDK design rules.",
                "Self-Healing: Synthesizes corrective offset constraints (e.g. shift x by 0.1um) and loops until zero DRC violations remain."
            ]
        },
        # Slide 9: Intro - Presentation Overview
        {
            "type": "split",
            "title": "Introduction: Presentation Overview",
            "box1_title": "Research Questions Addressed",
            "box1_bullets": [
                "Can LLM agents reason about physical analog layout placement topologies without coordinate hallucinations?",
                "How can we enforce sub-micron design rules and track snapping in a generative AI layout pipeline?",
                "Can we enable real-time interactive user inputs via conversational interfaces?"
            ],
            "box2_title": "Presentation Structure",
            "box2_bullets": [
                "System Architecture: Decoupled Strategic/Geometric paradigm.",
                "Stage 1 & 2: Ingestion, Graph Symmetry, and LangGraph Placer-Critic.",
                "Stage 3 & 4: PySide6 GUI, XML Command Queues, and Compaction Exporters.",
                "Evaluation: Real 14nm FinFET layouts, area footprints, and parasitics."
            ]
        },
        # Slide 10: Section 2 Title
        {
            "type": "section",
            "title": "Section 2: System Architecture",
            "subtitle": "Strategic Planning, Geometric Snapping, and the Multi-Agent Framework"
        },
        # Slide 11: System Architecture - 4 panels
        {
            "type": "four_panel"
        },
        # Slide 12: System Architecture - Decoupled Layers detail
        {
            "type": "split",
            "title": "System Architecture: Strategic & Geometric Separation",
            "box1_title": "Strategic Planning Layer (AI Agents)",
            "box1_bullets": [
                "Relative Ordering: Determines horizontal layout queues of cells (e.g. row0: MM0, MM1, MM2).",
                "Device Configurations: Recommends transistor finger splitting and unrolling to optimize aspect ratio.",
                "Symmetry Matching: Groups transistors into symmetrical common-centroid configurations.",
                "Abstraction: Completely isolated from sub-micron dimensions to eliminate spatial math errors."
            ],
            "box2_title": "Geometric Compiling Layer (Deterministic Engine)",
            "box2_bullets": [
                "Track Snapping: Aligns gate centers and active regions to quantized PDK fin and poly tracks.",
                "Diffusion Sharing: Merges adjacent transistor source/drains to reduce junction area.",
                "Clearance Compiler: Translates symbolic relative indices into physical spacing dimensions using PDK lookup rules."
            ]
        },
        # Slide 13: System Architecture - 4-Stage Pipeline detail
        {
            "type": "split",
            "title": "System Architecture: 4-Stage Compilation Pipeline",
            "box1_title": "Stages 1 & 2 (Backend Core)",
            "box1_bullets": [
                "Stage 1: EDA-Interface Inputs: Ingests SPICE netlists, resolves PDK parameters, and performs graph-isomorphism symmetry detection.",
                "Stage 2: AI-Initial Placement: Executes LangGraph cooperative state-machine planner and DRC self-healing critic loops."
            ],
            "box2_title": "Stages 3 & 4 (Frontend & Exporter)",
            "box2_bullets": [
                "Stage 3: ChatBot and GUI Integration: Manages human-in-the-loop conversational panels, intent routing, and thread-safe signal queues.",
                "Stage 4: EDA-Interface Outputs: Executes layout compaction, creates Gdstk OASIS files, and triggers OpenAccess WATCHDOG injection."
            ]
        },
        # Slide 14: System Architecture - Block Diagram (NATIVE BLOCK DIAGRAM!)
        {
            "type": "visual_pipeline"
        },
        # Slide 15: System Architecture - LangGraph State Machine (visual flowchart loop!)
        {
            "type": "visual_langgraph"
        },
        # Slide 16: System Architecture - Human in the Loop
        {
            "type": "split",
            "title": "System Architecture: Human-in-the-Loop Co-Pilot",
            "box1_title": "Bridging Conversational AI and Layout Canvas",
            "box1_bullets": [
                "Interactive Co-Design: AI generates initial layouts, and humans override details on the canvas.",
                "User Intent Parser: Chat panel processes natural language commands (e.g. \"align differential pair gates\").",
                "Asynchronous Bridge: Pushes commands to XML execution queues, updating the GUI without UI freezing."
            ],
            "box2_title": "Intent Translation",
            "box2_bullets": [
                "Commands are parsed into structured layouts.",
                "XML layouts are pushed to the compiler to snap boundaries.",
                "Tracks are updated instantly on the PySide6 canvas."
            ]
        },
        # Slide 17: System Architecture - State Variables
        {
            "type": "split",
            "title": "System Architecture: State Variables & Context",
            "box1_title": "LangGraph Active State Variables",
            "box1_bullets": [
                "netlist_graph: NetworkX representation of the electrical netlist.",
                "symmetry_groups: List of matched symmetry transistor groups.",
                "symbolic_map: JSON representation of active rows, devices, ordering, orientations, and fingers.",
                "coordinate_matrix: Absolute snapped micron coordinates."
            ],
            "box2_title": "Critic History Buffer",
            "box2_bullets": [
                "critic_history: Stores previous DRC check violations.",
                "correction_constraints: List of active spacing offset limits.",
                "user_overrides: Active list of designer override instructions."
            ]
        },
        # Slide 18: Section 3 Title
        {
            "type": "section",
            "title": "Section 3: Stage 1: EDA-Interface Inputs",
            "subtitle": "Netlist Parsing, Subgraph Isomorphism, and Parameter Resolution"
        },
        # Slide 19: Stage 1 - SPICE Ingestion
        {
            "type": "split",
            "title": "Stage 1: SPICE Ingestion & Parsing",
            "box1_title": "SPICE Netlist Parsing",
            "box1_bullets": [
                "Netlist Reader: Ingests flat or hierarchical SPICE files (.sp, .cdl).",
                "Subcircuit Extraction: Flattens modules and extracts individual transistor rows.",
                "Pin Mapping: Resolves electrical connections for transistor source, drain, gate, and body nodes."
            ],
            "box2_title": "Device Parameter Extraction",
            "box2_bullets": [
                "Transistor Dimensions: Resolves width (W), length (L), and finger counts (M).",
                "PDK Attributes: Maps discrete FinFET parameters such as fin count, gate pitch, and threshold options.",
                "Database Compilation: Builds a structured JSON node list of all devices."
            ]
        },
        # Slide 20: Stage 1 - Subgraph Isomorphism
        {
            "type": "split",
            "title": "Stage 1: Subgraph Isomorphism & Symmetry",
            "box1_title": "Graph Construction",
            "box1_bullets": [
                "Graph Representation: Ingested netlist is mapped to a bipartite graph using NetworkX.",
                "Bipartite Nodes: Devices and electrical nets are represented as distinct node groups.",
                "Bipartite Edges: Edges connect device terminal nodes to corresponding nets."
            ],
            "box2_title": "Isomorphism Matching",
            "box2_bullets": [
                "Pattern Recognition: Walks the bipartite graph checking for structural symmetry.",
                "Symmetry Matching: Identifies equivalent electrical configurations (e.g. symmetrical paths).",
                "Performance: Matching runs in milliseconds for typical analog blocks ($N < 50$ devices)."
            ]
        },
        # Slide 21: Stage 1 - Symmetry Groups
        {
            "type": "split",
            "title": "Stage 1: Symmetrical Matching Templates",
            "box1_title": "Differential Input Pairs",
            "box1_bullets": [
                "Symmetry Condition: Symmetrical transistors sharing gate inputs or cross-coupled gate-drain connections.",
                "Group Constraint: Forces transistors to be placed symmetrically around a central axis.",
                "Alignment: Aligns horizontal y-coordinates and snaps poly tracks."
            ],
            "box2_title": "Active Loads & Mirrors",
            "box2_bullets": [
                "Current Mirrors: Symmetrical transistor groups sharing common gate and source connections.",
                "Aspect Ratios: Identifies matching finger counts to ensure proportional current mirroring.",
                "Constraint Hierarchy: Writes symmetry requirements to the LangGraph state."
            ]
        },
        # Slide 22: Stage 1 - PDK track mapping
        {
            "type": "split",
            "title": "Stage 1: PDK Parameter Mapping & Snapping",
            "box1_title": "FinFET Grid Track Snapping",
            "box1_bullets": [
                "Ppoly Pitch: Gate centers must snap to integer multiples of the poly track pitch (e.g. 140nm).",
                "Pfin Pitch: Active diffusions must snap to horizontal fin track pitches (e.g. 48nm).",
                "Clearance Compiler: Reads spacing files and pre-calculates spacing margins."
            ],
            "box2_title": "Ingestion Boundaries",
            "box2_bullets": [
                "Cell Boundaries: Formulates minimum width and height boundaries for NMOS and PMOS regions.",
                "Well Clearances: Pre-calculates well isolation spacing limits to ensure DRC compliance."
            ]
        },
        # Slide 23: Stage 1 - Flowchart Diagram (NATIVE FLOWCHART!)
        {
            "type": "visual_parser"
        },
        # Slide 24: Stage 1 - Code Snippet
        {
            "type": "code",
            "title": "Stage 1: Parsed Ingestion Database Schema",
            "code_lines": [
                "{",
                "  \"nodes\": [",
                "    {",
                "      \"id\": \"MM0\",",
                "      \"type\": \"nmos\",",
                "      \"w\": 0.5e-6,",
                "      \"l\": 0.03e-6,",
                "      \"nf\": 2,",
                "      \"nets\": {\"d\": \"out\", \"g\": \"in\", \"s\": \"gnd\", \"b\": \"gnd\"}",
                "    },",
                "    {",
                "      \"id\": \"MM1\",",
                "      \"type\": \"nmos\",",
                "      \"w\": 0.5e-6,",
                "      \"l\": 0.03e-6,",
                "      \"nf\": 2,",
                "      \"nets\": {\"d\": \"out_b\", \"g\": \"in_b\", \"s\": \"gnd\", \"b\": \"gnd\"}",
                "    }",
                "  ],",
                "  \"edges\": [",
                "    {\"source\": \"MM0\", \"target\": \"MM1\", \"net\": \"gnd\", \"type\": \"s_s\"}",
                "  ]",
                "}"
            ]
        },
        # Slide 25: Section 4 Title
        {
            "type": "section",
            "title": "Section 4: Stage 2: AI-Initial Placement",
            "subtitle": "LangGraph Orchestration, Snapping, and Self-Healing Critic Loops"
        },
        # Slide 26: Stage 2 - Symbolic vs Physical
        {
            "type": "split",
            "title": "Stage 2: Symbolic vs. Physical Placement",
            "box1_title": "Eliminating Coordinate Hallucinations",
            "box1_bullets": [
                "The Problem: Generative AI models struggle with sub-micron coordinate arithmetic, causing DRC overlapping.",
                "The Solution: LLM agents are restricted to planning a symbolic placement map.",
                "No Math: LLM planning is restricted to relative ordering (e.g., placing M1 next to M2)."
            ],
            "box2_title": "Deterministic Snapping Compiler",
            "box2_bullets": [
                "Compile Snapping: A geometry compiler resolves symbolic indices into physical coordinates.",
                "Boundary Compiler: Integrates PDK boundary rules to calculate cell widths and heights."
            ]
        },
        # Slide 27: Stage 2 - Node Roles (visual grid!)
        {
            "type": "visual_nodes"
        },
        # Slide 28: Stage 2 - ReAct Self-Healing
        {
            "type": "split",
            "title": "Stage 2: Spacing Self-Healing Critic Loop",
            "box1_title": "Sweepline DRC Checks",
            "box1_bullets": [
                "Overlap Detection: Checks layout boundaries to identify overlaps.",
                "Clearance Verification: Measures gaps against PDK spacing rules."
            ],
            "box2_title": "Corrective Offset Synthesis",
            "box2_bullets": [
                "Constraint Injection: Generates relative spacing coordinates (e.g. shift x by 0.1um) when DRC violations are found.",
                "ReAct Loop: LangGraph loops from Specialist to Critic until all violations are resolved."
            ]
        },
        # Slide 29: Stage 2 - Flowchart Diagram (NATIVE FLOWCHART!)
        {
            "type": "visual_placement_flow"
        },
        # Slide 30: Stage 2 - Diffusion Sharing
        {
            "type": "split",
            "title": "Stage 2: Diffusion Sharing & Abutment",
            "box1_title": "Active diffusion sharing",
            "box1_bullets": [
                "Concept: Adjacent transistors sharing a common electrical net can share active diffusion regions.",
                "Contact Merging: Compactor merges source/drain contacts, eliminating isolation spacing.",
                "Area Benefits: Reduces cell width by up to 40%, optimizing layout density."
            ],
            "box2_title": "Abutment Engine",
            "box2_bullets": [
                "Abutment Search: Evaluates orientation flips (source/drain) to maximize shared contacts.",
                "Net Compatibility: Checks poly spacing and contact boundary rules."
            ]
        },
        # Slide 31: Stage 2 - Abutment & Dummies
        {
            "type": "split",
            "title": "Stage 2: Abutment Engine & Dummy Insertion",
            "box1_title": "Union-Find Chains",
            "box1_bullets": [
                "Abutment Chain Search: Identifies continuous transistor chains sharing source/drain nodes.",
                "Grid Snapping: Aligns active channels to horizontal tracks."
            ],
            "box2_title": "Conflict Spacing & Dummies",
            "box2_bullets": [
                "Abutment Conflicts: Inserts dummy isolation transistors or spacing guard zones if active nets conflict.",
                "PDK Snapping: Snaps boundaries to integer track multiples."
            ]
        },
        # Slide 32: Stage 2 - Code Snippet
        {
            "type": "code",
            "title": "Stage 2: Symbolic Placement JSON Output",
            "code_lines": [
                "{",
                "  \"rows\": [",
                "    {",
                "      \"id\": \"row0\", \"type\": \"nmos\", \"height\": 0.6,",
                "      \"devices\": [",
                "        {\"id\": \"MM0\", \"finger\": 0, \"orient\": \"R0\"},",
                "        {\"id\": \"MM0\", \"finger\": 1, \"orient\": \"MX\"},",
                "        {\"id\": \"MM1\", \"finger\": 0, \"orient\": \"R0\"},",
                "        {\"id\": \"MM1\", \"finger\": 1, \"orient\": \"MX\"}",
                "      ]",
                "    }",
                "  ]",
                "}"
            ]
        },
        # Slide 33: Section 5 Title
        {
            "type": "section",
            "title": "Section 5: Stage 3: ChatBot & GUI Integration",
            "subtitle": "Conversational Panel, Intent Classifiers, and Multi-Threaded Signal Slots"
        },
        # Slide 34: Stage 3 - Conversational UI
        {
            "type": "split",
            "title": "Stage 3: Conversational Layout Co-Pilot",
            "box1_title": "Interactive Chat Panel",
            "box1_bullets": [
                "Canvas Chat Panel: Integrated directly inside the PySide6 canvas editor.",
                "Natural Language Interface: Allows designers to interact with layout parameters via text.",
                "Query Support: Designers can query PDK rules, check symmetry constraints, or get coordinates."
            ],
            "box2_title": "Designer Overrides",
            "box2_bullets": [
                "Designer Command: Allows engineers to override placement settings.",
                "Intent Capture: Chat panel converts natural language commands into layout modifications."
            ]
        },
        # Slide 35: Stage 3 - Intent Classification (visual block diagram!)
        {
            "type": "visual_routing"
        },
        # Slide 36: Stage 3 - PySide6 GUI
        {
            "type": "split",
            "title": "Stage 3: PySide6 GUI Canvas Editor",
            "box1_title": "Interactive Layout Canvas",
            "box1_bullets": [
                "Canvas Editor: Displays transistor shapes (diffusions, gates, contacts, pins) dynamically.",
                "Graphic Items: Uses QGraphicsScene to display layout geometries."
            ],
            "box2_title": "Real-Time Updates",
            "box2_bullets": [
                "Real-Time Refreshes: Refreshes canvas layers instantly when placement updates are received.",
                "Visual Highlights: Highlights nets and symmetry axes."
            ]
        },
        # Slide 37: Stage 3 - Multi-Threaded Signal-Slot
        {
            "type": "split",
            "title": "Stage 3: Thread-Safe GUI Orchestration",
            "box1_title": "Thread Separation",
            "box1_bullets": [
                "OrchestratorWorker: Moves LLM requests and LangGraph iterations to a background thread.",
                "Thread Separation: Protects PySide6 GUI main thread from freezing during long reasoning runs."
            ],
            "box2_title": "Qt Signal-Slot Communication",
            "box2_bullets": [
                "Communication uses thread-safe signals:",
                "stage_started: Shows loading spinners in the GUI.",
                "response_delta: Streams text token-by-token to the chat panel.",
                "command_ready: Transmits layout updates."
            ]
        },
        # Slide 38: Stage 3 - Command Queue
        {
            "type": "split",
            "title": "Stage 3: Command Queue & Interactive Undo",
            "box1_title": "XML Layout Queues",
            "box1_bullets": [
                "XML Command Queue: Pushes all updates to a structured queue.",
                "Batch Command Flusher: Group commands into a single canvas redraw transaction using a 0ms QTimer to prevent flickering."
            ],
            "box2_title": "Undo/Redo Stack",
            "box2_bullets": [
                "Undo Stack: Tracks all state changes to allow reverting layout adjustments.",
                "Co-Design Support: Allows designers to easily undo or redo modifications."
            ]
        },
        # Slide 39: Stage 3 - Flowchart Diagram (NATIVE FLOWCHART!)
        {
            "type": "visual_chatbot_flow"
        },
        # Slide 40: Section 6 Title
        {
            "type": "section",
            "title": "Section 6: Stage 4: EDA-Interface Outputs",
            "subtitle": "Compaction, OASIS Exporter, and OpenAccess TCL Injection"
        },
        # Slide 41: Stage 4 - Physical Compaction
        {
            "type": "split",
            "title": "Stage 4: Physical Layout Compaction",
            "box1_title": "Compaction Execution",
            "box1_bullets": [
                "Compaction Engine: Ingests placement maps and resolves PDK spacing rules.",
                "Clearance Sweeper: Compacts gaps between gates and active areas to minimize overall cell area.",
                "Pin Alignment: Snaps contacts and pins to vertical and horizontal routing grids."
            ],
            "box2_title": "LDE Reduction",
            "box2_bullets": [
                "WPE and STI Protection: Places dummy devices to shield active transistors.",
                "Symmetry Axis Centering: Aligns symmetrical transistor rows."
            ]
        },
        # Slide 42: Stage 4 - OASIS Exporter
        {
            "type": "split",
            "title": "Stage 4: Binary GDSII & OASIS Exporter",
            "box1_title": "Gdstk Layout Assembly",
            "box1_bullets": [
                "Gdstk Exporter: Generates layouts using the C++/Python Gdstk library.",
                "Cell Assembly: Clones PDK base PCells and compiles geometry polygons.",
                "Speed: Layout generation finishes in milliseconds."
            ],
            "box2_title": "Binary Formats",
            "box2_bullets": [
                "OASIS Exporter: Exports binary GDSII or OASIS layout streams.",
                "OASIS Advantage: Offers up to 10x smaller file sizes than GDSII, accelerating large macro exports."
            ]
        },
        # Slide 43: Stage 4 - OpenAccess watchdog (NATIVE FLOWCHART!)
        {
            "type": "visual_watchdog"
        },
        # Slide 44: Stage 4 - Exporter Flowchart Diagram (NATIVE FLOWCHART!)
        {
            "type": "visual_exporter_flow"
        },
        # Slide 45: Stage 4 - Code Snippet
        {
            "type": "code",
            "title": "Stage 4: Watchdog Observer TCL Script",
            "code_lines": [
                "# cc_watcher.tcl - Ingest JSON layout updates",
                "if {[file exists $layout_json_path]} {",
                "    set fp [open $layout_json_path r]",
                "    set file_data [read $fp]",
                "    close $fp",
                "    # Parse placement coordinate arrays",
                "    set json_data [json::json2dict $file_data]",
                "    foreach dev_id [dict keys $json_data] {",
                "        set coords [dict get $json_data $dev_id]",
                "        set xPos [lindex $coords 0]",
                "        set yPos [lindex $coords 1]",
                "        # Inject coordinates directly into active cellview",
                "        db_move_device $dev_id $xPos $yPos",
                "    }",
                "}"
            ]
        },
        # Slide 46: Section 7 Title
        {
            "type": "section",
            "title": "Section 7: Evaluation & Results",
            "subtitle": "XOR, Mirror, and Comparator Layout Benchmarks, Area, Parasitics, and Latency"
        },
        # Slide 47: Evaluation - Experimental Setup
        {
            "type": "split",
            "title": "Evaluation: Experimental Setup",
            "box1_title": "Process Design Kits & Circuits",
            "box1_bullets": [
                "ASU 14nm: Predictive FinFET PDK.",
                "SAED 14nm: Synopsys FinFET PDK.",
                "Circuits: 2-Transistor Mirror, Differential Input Pair, 5-T OTA, Dynamic Comparator, XOR Gate."
            ],
            "box2_title": "Compute Specs & Models",
            "box2_bullets": [
                "Compute Spec: Intel Core i9-12900K, 64GB RAM, RTX 3090 (24GB VRAM).",
                "LLM Backends: Claude 3.5 Sonnet, GPT-4o, Llama-3-8B (Ollama Local).",
                "EDA Tools: Synopsys Custom Compiler, Cadence Virtuoso, StarRC."
            ]
        },
        # Slide 48: Evaluation - XOR Gate placement images (Max 2 images!)
        {
            "type": "image_grid_comparator",
            "title": "Evaluation: XOR Gate Layout - Symbolic Placements",
            "img_a": "xor_before_AI_placement.png",
            "img_b": "xor_symbolic.png"
        },
        # Slide 49: Evaluation - XOR Gate physical images (Max 2 images!)
        {
            "type": "image_grid_comparator",
            "title": "Evaluation: XOR Gate Layout - Schematic & Compacted",
            "img_a": "xor_transistor_level.png",
            "img_b": "xor_klayout.png"
        },
        # Slide 50: Evaluation - Current Mirror placement (Max 2 images!)
        {
            "type": "image_grid_comparator",
            "title": "Evaluation: Current Mirror Layout - Symbolic Placements",
            "img_a": "current_mirror_before_AI_placement.png",
            "img_b": "current_mirror_symbolic.png"
        },
        # Slide 51: Evaluation - Current Mirror schematic and centroid (Max 2 images!)
        {
            "type": "image_grid_comparator",
            "title": "Evaluation: Current Mirror Layout - Schematic & centroid Matching",
            "img_a": "current_mirror_transistor_level.png",
            "img_b": "current_mirror_transistor_level_colord.png"
        },
        # Slide 52: Evaluation - Current Mirror KLayout and imported (Max 2 images!)
        {
            "type": "image_grid_comparator",
            "title": "Evaluation: Current Mirror Layout - Compacted & EDA Ingestion",
            "img_a": "current_mirror_klayout.png",
            "img_b": "current_mirror_from_custom_compiler.png"
        },
        # Slide 53: Evaluation - Dynamic Comparator 1x2 Image Grid (Max 2 images!)
        {
            "type": "image_grid_comparator",
            "title": "Evaluation: Dynamic Comparator Layout",
            "img_a": "comparator_before_AI_placement.png",
            "img_b": "Comparator.png"
        },
        # Slide 54: Evaluation - Area Reductions Table
        {
            "type": "table_only",
            "title": "Evaluation: Layout Area Footprints",
            "headers": ["Circuit Block", "Symbolic Area (um2)", "Compacted Area (um2)", "Manual Area (um2)", "Area Savings", "Gap to Manual"],
            "rows": [
                ["2-T Mirror", "0.588", "0.350", "0.345", "40.5%", "+1.4%"],
                ["Diff. Pair", "1.176", "0.770", "0.750", "34.5%", "+2.6%"],
                ["5-T OTA", "2.940", "1.820", "1.780", "38.1%", "+2.2%"],
                ["Comparator", "5.880", "3.710", "3.550", "36.9%", "+4.5%"],
                ["XOR Gate", "1.470", "0.920", "0.900", "37.4%", "+2.2%"]
            ]
        },
        # Slide 55: Evaluation - Parasitics Table
        {
            "type": "table_only",
            "title": "Evaluation: Post-Layout Parasitics",
            "headers": ["Circuit Block", "Isolated Cdb,sb (fF)", "Compacted Cdb,sb (fF)", "Junction Savings", "Wire R (Ohm)", "Wire C (fF)"],
            "rows": [
                ["2-T Mirror", "0.88", "0.22", "75.0%", "12.4", "0.15"],
                ["Diff. Pair", "1.76", "0.56", "68.2%", "24.8", "0.32"],
                ["5-T OTA", "4.40", "1.23", "72.0%", "48.5", "0.65"],
                ["Comparator", "8.80", "2.60", "70.5%", "92.4", "1.28"],
                ["XOR Gate", "2.20", "0.68", "69.1%", "35.6", "0.44"]
            ]
        },
        # Slide 56: Evaluation - Latency Table
        {
            "type": "table_only",
            "title": "Evaluation: Conversational Execution Latency",
            "headers": ["LLM Backend", "Inference (ms)", "Orchestration (ms)", "GUI & TCP (ms)", "Total Loop (ms)"],
            "rows": [
                ["GPT-4o (Cloud)", "420", "15", "12", "447"],
                ["Claude 3.5 (Cloud)", "580", "15", "12", "607"],
                ["Llama-3-8B (Local)", "210", "14", "12", "236"]
            ]
        },
        # Slide 57: Evaluation - Self-Healing Convergence Table
        {
            "type": "table_only",
            "title": "Evaluation: LangGraph Self-Healing Convergence",
            "headers": ["Iteration", "Current Mirror", "Diff Input Pair", "5-Transistor OTA", "Dynamic Comparator"],
            "rows": [
                ["0 (Initial Plan)", "3 violations", "5 violations", "12 violations", "24 violations"],
                ["1", "0 violations", "2 violations", "4 violations", "8 violations"],
                ["2", "0 violations", "0 violations", "1 violation", "3 violations"],
                ["3", "0 violations", "0 violations", "0 violations", "1 violation"],
                ["4 (DRC Clean)", "0 violations", "0 violations", "0 violations", "0 violations"]
            ]
        },
        # Slide 58: Section 8 Title
        {
            "type": "section",
            "title": "Section 8: Conclusions & Future Work",
            "subtitle": "Core Contributions, Technical Challenges, and AI Routing Roadmap"
        },
        # Slide 59: Future Work & Conclusion (rendered as multi-card)
        {
            "type": "future_work"
        },
        # Slide 60: Final Thank You
        {
            "type": "thank_you",
            "title": "Thank You!",
            "subtitle": "AI-Based Analog Layout Automation: A Conversational, Multi-Agent Framework",
            "author": "Joint Collaboration: Ain Shams University & SI-VISION",
            "institute": "Questions & Answers",
            "date": "June 2026",
            "logo_aus": "ASU.png",
            "logo_si": "SI.png"
        }
    ]
    
    total_pages = len(slides_data)
    
    for idx, data in enumerate(slides_data, 1):
        slide = prs.slides.add_slide(blank_layout)
        
        if data["type"] == "cover":
            add_cover_slide(slide, data)
            
        elif data["type"] == "timeline":
            add_timeline_slide(slide, idx, total_pages)
            
        elif data["type"] == "section":
            # Solid Academic Blue background
            add_solid_background(slide, ACADEMIC_BLUE)
            
            # Center section title
            tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(1.5))
            tf = tx_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = data["title"]
            p.font.name = "Calibri"
            p.font.size = Pt(38)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
            # Accent bar
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), Inches(3.8), Inches(5.333), Inches(0.05))
            line.fill.solid()
            line.fill.fore_color.rgb = WARM_GOLD
            line.line.fill.background()
            
            # Subtitle
            tx_box2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.1), Inches(11.333), Inches(1.2))
            tf2 = tx_box2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = data["subtitle"]
            p2.font.name = "Calibri"
            p2.font.size = Pt(18)
            p2.font.color.rgb = WARM_GOLD
            p2.alignment = PP_ALIGN.CENTER
            
            # Add slide footer on dark blue background
            add_slide_footer(slide, idx, total_pages, is_dark=True)
            
        elif data["type"] == "revolution":
            add_revolution_slide(slide, idx, total_pages)
            
        elif data["type"] == "split":
            add_slide_header(slide, data["title"])
            add_slide_footer(slide, idx, total_pages)
            
            # Left & Right content boxes
            box_width = Inches(5.6)
            box_height = Inches(5.4)
            add_content_box(slide, Inches(0.8), Inches(1.4), box_width, box_height, data["box1_title"], data["box1_bullets"])
            add_content_box(slide, Inches(6.8), Inches(1.4), box_width, box_height, data["box2_title"], data["box2_bullets"])
            
        elif data["type"] == "image_split":
            add_slide_header(slide, data["title"])
            add_slide_footer(slide, idx, total_pages)
            
            # Left content box, Right boxed image
            add_content_box(slide, Inches(0.8), Inches(1.4), Inches(5.2), Inches(5.4), data["box_title"], data["bullets"])
            add_boxed_image(slide, Inches(6.6), Inches(1.4), Inches(5.8), Inches(5.0), data["image"], data["label"])
            
        elif data["type"] == "four_panel":
            add_four_panel_slide(slide, idx, total_pages)
            
        elif data["type"] == "visual_pipeline":
            add_visual_pipeline(slide, idx, total_pages)
            
        elif data["type"] == "visual_langgraph":
            add_visual_langgraph(slide, idx, total_pages)
            
        elif data["type"] == "visual_nodes":
            add_visual_nodes(slide, idx, total_pages)
            
        elif data["type"] == "visual_routing":
            add_visual_routing(slide, idx, total_pages)
            
        elif data["type"] == "visual_watchdog":
            add_visual_watchdog(slide, idx, total_pages)
            
        elif data["type"] == "visual_parser":
            add_visual_parser(slide, idx, total_pages)
            
        elif data["type"] == "visual_placement_flow":
            add_visual_placement_flow(slide, idx, total_pages)
            
        elif data["type"] == "visual_chatbot_flow":
            add_visual_chatbot_flow(slide, idx, total_pages)
            
        elif data["type"] == "visual_exporter_flow":
            add_visual_exporter_flow(slide, idx, total_pages)
            
        elif data["type"] == "code":
            add_slide_header(slide, data["title"])
            add_slide_footer(slide, idx, total_pages)
            
            # Decode explanations and list side by side
            if "Database Schema" in data["title"]:
                explanations = [
                    "Nodes list: Represents devices in the flattened netlist (e.g. MM0, MM1).",
                    "Pins mapping: Links gate (g), source (s), drain (d), and bulk (b) terminals to logical net names.",
                    "Properties fields: Maps physical FinFET parameters: width (w), length (l), and parallel split fingers (nf)."
                ]
                box_title = "JSON Ingestion Fields"
                code_title = "Parsed Netlist JSON Schema"
            elif "Watchdog Observer" in data["title"]:
                explanations = [
                    "File Observer: Monitors the filesystem directory for updates to the coordinate exchange file.",
                    "JSON Parsing: Ingests layout changes as key-value pairs (device ID to physical coordinates).",
                    "Database Injection: Executes db_move_device to update the active OpenAccess cellview in real-time."
                ]
                box_title = "TCL Script Logic"
                code_title = "cc_watcher.tcl Watchdog Script"
            else:
                explanations = [
                    "Rows template: Strategic row partitions separating NMOS and PMOS device regions.",
                    "Relative index: Horizontal sequence numbers representing device ordering from left to right.",
                    "Orientations flags: Flip states (R0 vs. MX) that control source/drain terminal configurations to allow contact sharing."
                ]
                box_title = "JSON Symbolic Fields"
                code_title = "Symbolic Placement JSON Output"
                
            add_content_box(slide, Inches(0.8), Inches(1.4), Inches(5.0), Inches(5.4), box_title, explanations)
            add_code_box(slide, Inches(6.3), Inches(1.4), Inches(6.2), Inches(5.4), code_title, data["code_lines"])
            
        elif data["type"] == "table_only":
            add_slide_header(slide, data["title"])
            add_slide_footer(slide, idx, total_pages)
            
            # Centered table
            add_table_box(slide, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.5), data["headers"], data["rows"])
            
        elif data["type"] == "image_grid_xor":
            # (Backup block in case used, but split XOR is used instead)
            add_slide_header(slide, data["title"])
            add_slide_footer(slide, idx, total_pages)
            add_boxed_image(slide, Inches(1.2), Inches(1.5), Inches(5.2), Inches(2.2), data["img_a"], "(a) Initial unplaced symbolic grid.")
            add_boxed_image(slide, Inches(6.8), Inches(1.5), Inches(5.2), Inches(2.2), data["img_b"], "(b) Symbolic placement with abutment.")
            add_boxed_image(slide, Inches(1.2), Inches(4.3), Inches(5.2), Inches(2.2), data["img_c"], "(c) Transistor schematic.")
            add_boxed_image(slide, Inches(6.8), Inches(4.3), Inches(5.2), Inches(2.2), data["img_d"], "(d) Final compacted GDSII cell.")
            
        elif data["type"] == "image_grid_mirror":
            # (Backup block in case used, but split CM is used instead)
            add_slide_header(slide, data["title"])
            add_slide_footer(slide, idx, total_pages)
            add_boxed_image(slide, Inches(0.8), Inches(1.5), Inches(3.6), Inches(1.8), data["img_a"], "(a) Initial uncompacted placement.")
            add_boxed_image(slide, Inches(4.8), Inches(1.5), Inches(3.6), Inches(1.8), data["img_b"], "(b) Abutted placement.")
            add_boxed_image(slide, Inches(8.8), Inches(1.5), Inches(3.6), Inches(1.8), data["img_c"], "(c) Transistor schematic.")
            add_boxed_image(slide, Inches(0.8), Inches(4.2), Inches(3.6), Inches(1.8), data["img_d"], "(d) Centroid plan.")
            add_boxed_image(slide, Inches(4.8), Inches(4.2), Inches(3.6), Inches(1.8), data["img_e"], "(e) GDSII layout.")
            add_boxed_image(slide, Inches(8.8), Inches(4.2), Inches(3.6), Inches(1.8), data["img_f"], "(f) Virtuoso view.")
            
        elif data["type"] == "image_grid_comparator":
            add_slide_header(slide, data["title"])
            add_slide_footer(slide, idx, total_pages)
            
            # Resolve image labels based on slide title
            if "XOR" in data["title"]:
                if "Symbolic Placements" in data["title"]:
                    label_a = "(a) Initial unplaced symbolic grid."
                    label_b = "(b) Symbolic placement with abutment."
                else:
                    label_a = "(c) Transistor schematic."
                    label_b = "(d) Final compacted GDSII cell in KLayout."
            elif "Current Mirror" in data["title"]:
                if "Symbolic Placements" in data["title"]:
                    label_a = "(a) Initial uncompacted placement."
                    label_b = "(b) Abutted symbolic placement."
                elif "Schematic" in data["title"]:
                    label_a = "(c) Transistor schematic."
                    label_b = "(d) Color-coded common-centroid plan."
                else:
                    label_a = "(e) Compacted GDSII in KLayout."
                    label_b = "(f) Imported OpenAccess cell view."
            else:
                label_a = "(a) Initial unplaced schematic cells."
                label_b = "(b) Final physical layout with snapping in KLayout."
                
            # Render two large side-by-side images centered on the slide (no analysis panel)
            add_boxed_image(slide, Inches(0.8), Inches(1.5), Inches(5.4), Inches(4.5), data["img_a"], label_a)
            add_boxed_image(slide, Inches(7.1), Inches(1.5), Inches(5.4), Inches(4.5), data["img_b"], label_b)
            
        elif data["type"] == "future_work":
            add_future_work_slide(slide, idx, total_pages)
            
        elif data["type"] == "thank_you":
            # Solid Academic Blue background
            add_solid_background(slide, ACADEMIC_BLUE)
            
            # Title
            tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(1.8))
            tf = tx_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = data["title"]
            p.font.name = "Calibri"
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
            # Gold Divider line
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.0), Inches(3.5), Inches(7.333), Inches(0.04))
            line.fill.solid()
            line.fill.fore_color.rgb = WARM_GOLD
            line.line.fill.background()
            
            # Subtitle
            tx_box2 = slide.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(11.333), Inches(1.0))
            tf2 = tx_box2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = data["subtitle"]
            p2.font.name = "Calibri"
            p2.font.size = Pt(18)
            p2.font.color.rgb = WARM_GOLD
            p2.alignment = PP_ALIGN.CENTER
            
            # Presenter names
            tx_box3 = slide.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.333), Inches(1.0))
            tf3 = tx_box3.text_frame
            p3 = tf3.paragraphs[0]
            p3.text = data["author"]
            p3.font.name = "Calibri"
            p3.font.size = Pt(14)
            p3.font.bold = True
            p3.font.color.rgb = WHITE
            p3.alignment = PP_ALIGN.CENTER
            
            # Institute
            p_inst = tf3.add_paragraph()
            p_inst.text = data["institute"]
            p_inst.font.name = "Calibri"
            p_inst.font.size = Pt(12)
            p_inst.font.color.rgb = LIGHT_BLUE_BG
            p_inst.space_before = Pt(4)
            p_inst.alignment = PP_ALIGN.CENTER
            
            # Date
            p_date = tf3.add_paragraph()
            p_date.text = data["date"]
            p_date.font.name = "Calibri"
            p_date.font.size = Pt(10)
            p_date.font.color.rgb = WARM_GOLD
            p_date.space_before = Pt(4)
            p_date.alignment = PP_ALIGN.CENTER
            
            # Add logos side-by-side
            logo_aus_path = get_image_path(data["logo_aus"])
            logo_si_path = get_image_path(data["logo_si"])
            
            if logo_aus_path:
                slide.shapes.add_picture(logo_aus_path, Inches(1.5), Inches(5.8), height=Inches(1.0))
            if logo_si_path:
                slide.shapes.add_picture(logo_si_path, Inches(10.2), Inches(5.8), height=Inches(0.9))
                
    # Save Presentation
    prs.save("presentation.pptx")
    print("PowerPoint presentation generated successfully!")

if __name__ == "__main__":
    create_presentation()
